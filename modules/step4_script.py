# modeling/modules/step4_script.py

import os
import json
import sys
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv

# 상위 폴더(utils)를 찾기 위한 경로 설정
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# [경로 수정] 통합된 Gemini 클라이언트 불러오기
from utils.gemini_client import generate_script_and_qna

load_dotenv()

def extract_text_from_pptx(pptx_path):
    """PPT 파일에서 텍스트 추출"""
    if not os.path.exists(pptx_path):
        return None
    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text.strip() if slide.shapes.title else "무제"
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                if shape == slide.shapes.title: continue
                for p in shape.text_frame.paragraphs:
                    if p.text.strip(): content.append(p.text.strip())
        slides_text.append(f"[[Slide {i+1}]] Title: {title}\nContent: {' '.join(content)}")
    return "\n\n".join(slides_text)

def main(): # 외부 호출용 함수명
    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    input_path = os.path.join(BASE_DIR, "data", "script_input", "ppt_ex.pptx")
    
    print("="*60)
    print("[Step 4] PPT 대본 및 Q&A 생성")
    print("="*60)

    # 1. PPT 텍스트 추출
    print(f"[*] PPT 파일 읽는 중: {os.path.basename(input_path)}")
    ppt_text = extract_text_from_pptx(input_path)
    
    if not ppt_text: 
        print(f"[!] PPT 파일을 찾을 수 없습니다: {input_path}")
        return

    # 2. Gemini 호출 (utils/gemini_client.py 사용)
    # 복잡한 프롬프트 코드는 이제 gemini_client 안에 있습니다.
    json_data = generate_script_and_qna(ppt_text)
    
    # 3. 결과 저장
    if json_data:
        output_path = os.path.join(BASE_DIR, "data", "report", "script_flow.json")
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
            
        print(f"\n[*] 실전 발표용 데이터 생성 완료!")
        print(f"    - 구조: 슬라이드별 대본 + 마지막 종합 Q&A")
        print(f"    - 파일 위치: {os.path.abspath(output_path)}")

    print("="*60)
    print("Step 4 완료!")

if __name__ == "__main__":
    run_step4()