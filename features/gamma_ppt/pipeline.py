"""
Pipeline: PDF -> Gemini -> Gamma -> PPTX.
"""

from __future__ import annotations

import json
import os
import re
from datetime import datetime
from typing import Any, Dict, List, Optional

from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation

from utils.document_parsing import extract_text_from_pdf

from .gamma_client import (
    build_gamma_payload,
    build_template_payload,
    choose_pptx_url,
    create_generation,
    create_generation_from_template,
    download_pptx,
    poll_generation,
)
from .prompts import (
    SYSTEM_INSTRUCTION_CARDS,
    SYSTEM_INSTRUCTION_CONDENSE,
    SYSTEM_INSTRUCTION_DECK,
    SYSTEM_INSTRUCTION_TEMPLATE_JSON,
    build_condense_prompt,
    build_deck_prompt,
    build_recommended_cards_prompt,
    build_template_json_prompt,
)

load_dotenv()

GEMINI_MODEL_NAME = "gemini-2.5-flash"

MAX_RFP_CHARS = 100000
CONDENSE_TARGET_CHARS = 50000
CONDENSE_TARGET_CHARS_TEMPLATE = 35000
CONDENSE_TARGET_CHARS_FAST = 30000
CONDENSE_TARGET_CHARS_TEMPLATE_FAST = 25000
MIN_CONDENSE_CHARS = 20000

DEFAULT_CARDS = 18
MIN_CARDS = 16
MAX_CARDS = 28
TEMPLATE_SLIDE_COUNT = 18

MIN_BULLETS = 3
MAX_BULLETS = 5
MAX_BULLETS_IMAGE = 4
MAX_BULLET_CHARS = 140
SPLIT_CHAR_LIMIT = 900
SPLIT_CHAR_LIMIT_IMAGE = 650
MERGE_CHAR_LIMIT = 240

BANNED_PHRASES = [
    "차트", "그래프", "표로 표시", "표로", "그래프로",
    "bar chart", "todo", "tbd", "확인 필요",
]

IMAGE_ALLOWED_KEYWORDS = [
    "표지",
    "개요",
    "아키텍처",
    "구성",
    "로드맵",
    "추진체계",
    "기대효과",
    "기술",
    "시스템",
    "모듈",
    "데이터",
    "파이프라인",
    "프로세스",
    "워크플로우",
    "모델",
    "프레임워크",
    "플랫폼",
]

IMAGE_ERROR_PHRASES = [
    "이미지를 생성하는 중 오류",
    "이미지 생성 오류",
    "Error generating image",
    "Image generation failed",
]

DEFAULT_IMAGE_STYLE = "technical diagram, schematic, flat, monochrome, minimal, line art, blue-gray"

BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
DATA_DIR = os.path.join(BASE_DIR, "data")
EXPORT_DIR = os.path.join(DATA_DIR, "gamma_exports")
RESP_DIR = os.path.join(DATA_DIR, "gamma_responses")
REPORT_DIR = os.path.join(DATA_DIR, "gamma_reports")


def _gemini_generate(system_instruction: str, prompt: str, mime_type: str, temperature: float) -> str:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set.")

    client = genai.Client(api_key=api_key)
    response = client.models.generate_content(
        model=GEMINI_MODEL_NAME,
        contents=prompt,
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            response_mime_type=mime_type,
            temperature=temperature,
        ),
    )
    return response.text.strip()


def parse_pdf_to_text(pdf_path: str) -> str:
    pages = extract_text_from_pdf(pdf_path)
    chunks: List[str] = []
    for page in pages:
        page_index = page.get("page_index", 0)
        texts = page.get("texts", [])
        chunks.append(f"[Page {page_index + 1}]")
        chunks.extend(texts)
    return "\n".join(chunks)


def _condense_if_needed(rfp_text: str, target_chars: int) -> str:
    if len(rfp_text) <= target_chars:
        return rfp_text

    truncated = rfp_text[:MAX_RFP_CHARS]
    prompt = build_condense_prompt(truncated, target_chars)
    condensed = _gemini_generate(SYSTEM_INSTRUCTION_CONDENSE, prompt, "text/plain", 0.1)
    return condensed


def _compute_condense_target(raw_len: int, mode: str, condense_auto: bool) -> int:
    base = CONDENSE_TARGET_CHARS_TEMPLATE if mode == "template" else CONDENSE_TARGET_CHARS
    if not condense_auto:
        return base
    auto = int(raw_len * 0.35)
    auto = max(MIN_CONDENSE_CHARS, min(auto, base))
    cap = CONDENSE_TARGET_CHARS_TEMPLATE_FAST if mode == "template" else CONDENSE_TARGET_CHARS_FAST
    return min(auto, cap)


def _parse_recommended_cards(text: str) -> Optional[int]:
    try:
        data = json.loads(text)
        value = data.get("recommended_cards")
        if isinstance(value, int):
            return value
        if isinstance(value, str) and value.isdigit():
            return int(value)
    except Exception:
        pass

    m = re.search(r"(\d{1,2})", text)
    if m:
        try:
            return int(m.group(1))
        except Exception:
            return None
    return None


def _recommend_card_count(rfp_text: str, min_cards: int, max_cards: int) -> int:
    prompt = build_recommended_cards_prompt(rfp_text, min_cards, max_cards)
    text = _gemini_generate(SYSTEM_INSTRUCTION_CARDS, prompt, "application/json", 0.1)
    n = _parse_recommended_cards(text)
    if n is None:
        return DEFAULT_CARDS
    return n


def _clamp_cards(value: int, min_cards: int, max_cards: int) -> int:
    return max(min_cards, min(max_cards, value))


def _generate_deck_text(rfp_text: str, num_cards: int) -> str:
    prompt = build_deck_prompt(rfp_text, num_cards)
    return _gemini_generate(SYSTEM_INSTRUCTION_DECK, prompt, "text/plain", 0.2)


def _generate_template_json(rfp_text: str) -> str:
    prompt = build_template_json_prompt(rfp_text)
    return _gemini_generate(SYSTEM_INSTRUCTION_TEMPLATE_JSON, prompt, "application/json", 0.2)


def _parse_template_slides(text: str) -> List[Dict[str, Any]]:
    data = json.loads(text)
    slides = data.get("slides", [])
    if not isinstance(slides, list):
        raise ValueError("Invalid template JSON: slides is not a list.")
    normalized: List[Dict[str, Any]] = []
    for s in slides:
        if not isinstance(s, dict):
            continue
        title = str(s.get("title", "")).strip()
        bullets = s.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        if not isinstance(bullets, list):
            bullets = []
        bullets = [f"- {str(b).strip()}" if not str(b).strip().startswith("-") else str(b).strip() for b in bullets if str(b).strip()]
        normalized.append({"title": title, "bullets": bullets})
    return normalized


def _slides_to_deck_text(slides: List[Dict[str, Any]]) -> str:
    blocks: List[str] = []
    for idx, slide in enumerate(slides, 1):
        title = slide.get("title") or f"슬라이드 {idx}"
        bullets = slide.get("bullets") or []
        if not bullets:
            bullets = ["- 자료 없음"]
        blocks.append("\n".join([f"# {title}"] + bullets))
    return "\n---\n".join(blocks)


def _is_banned_line(line: str) -> bool:
    low = line.lower()
    for kw in BANNED_PHRASES:
        if kw in low:
            return True
    return False


def _title_has_image_keyword(title: str) -> bool:
    low = title.lower()
    for kw in IMAGE_ALLOWED_KEYWORDS:
        if kw.lower() in low:
            return True
    return False


def _max_bullets_for_title(title: str) -> int:
    return MAX_BULLETS_IMAGE if _title_has_image_keyword(title) else MAX_BULLETS


def _trim_bullet(bullet: str) -> str:
    text = bullet.strip()
    if text.startswith("-"):
        text = text.lstrip("-").strip()
    if len(text) <= MAX_BULLET_CHARS:
        return f"- {text}"
    cut = max(10, MAX_BULLET_CHARS - 3)
    return f"- {text[:cut].rstrip()}..."


def _clean_deck_text(deck_text: str) -> str:
    blocks = [b.strip() for b in deck_text.split("\n---\n")]
    cleaned_blocks: List[str] = []

    for idx, block in enumerate(blocks, 1):
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        title = None
        bullets: List[str] = []

        for ln in lines:
            if ln.startswith("#") and title is None:
                title = ln.lstrip("#").strip()
                continue
            if ln.startswith("-"):
                if not _is_banned_line(ln):
                    bullets.append(ln)

        if not title:
            title = f"슬라이드 {idx}"

        if not bullets:
            bullets = ["- 관련 자료 없음"]

        bullets = [_trim_bullet(b) for b in bullets]
        cleaned_blocks.append({"title": title, "bullets": bullets})

    adjusted = _balance_density(cleaned_blocks)

    final_blocks: List[str] = []
    for idx, block in enumerate(adjusted, 1):
        title = block["title"] or f"슬라이드 {idx}"
        bullets = block["bullets"]

        if not bullets:
            bullets = ["- 관련 자료 없음"]

        while len(bullets) < MIN_BULLETS:
            bullets.append("- 관련 자료 없음")

        max_bullets = _max_bullets_for_title(title)
        if len(bullets) > max_bullets:
            bullets = bullets[:max_bullets]

        final_blocks.append("\n".join([f"# {title}"] + bullets))

    return "\n---\n".join(final_blocks)


def _bullet_char_count(bullets: List[str]) -> int:
    return sum(len(b.replace("-", "").strip()) for b in bullets)


def _balance_density(blocks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    i = 0
    while i < len(blocks):
        block = blocks[i]
        title = block.get("title", "")
        bullets = block.get("bullets", [])
        bullet_count = len(bullets)
        char_count = _bullet_char_count(bullets)
        max_bullets = _max_bullets_for_title(title)
        split_limit = SPLIT_CHAR_LIMIT_IMAGE if _title_has_image_keyword(title) else SPLIT_CHAR_LIMIT

        # Split over-dense
        if bullet_count > max_bullets or char_count > split_limit:
            mid = max(1, min(len(bullets) - 1, len(bullets) // 2))
            first = bullets[:mid]
            second = bullets[mid:]
            out.append({"title": f"{title} (1/2)", "bullets": first})
            out.append({"title": f"{title} (2/2)", "bullets": second})
            i += 1
            continue

        # Merge under-dense with next if possible
        if bullet_count < MIN_BULLETS and char_count < MERGE_CHAR_LIMIT and i + 1 < len(blocks):
            next_block = blocks[i + 1]
            next_bullets = next_block.get("bullets", [])
            combined = bullets + next_bullets
            if _bullet_char_count(combined) <= SPLIT_CHAR_LIMIT:
                combined_title = f"{title} / {next_block.get('title', '')}".strip(" /")
                out.append({"title": combined_title, "bullets": combined})
                i += 2
                continue

        out.append({"title": title, "bullets": bullets})
        i += 1

    return out


def _get_image_error_phrases() -> List[str]:
    phrases = list(IMAGE_ERROR_PHRASES)
    extra = os.environ.get("GAMMA_IMAGE_ERROR_PHRASES")
    if extra:
        phrases.extend([p.strip() for p in extra.split(",") if p.strip()])
    return phrases


def _image_policy_instruction() -> str:
    keywords = ", ".join(IMAGE_ALLOWED_KEYWORDS)
    return (
        "이미지는 다음 키워드가 포함된 카드에만 사용하세요: "
        f"{keywords}. "
        "그 외 카드에는 이미지를 사용하지 마세요. "
        "이미지는 기술 설명용 도식/플로우차트/아키텍처 블록 다이어그램 위주로 생성하고 "
        "이미지 내 텍스트는 1~3개 이내의 아주 짧은 한글 라벨만 허용합니다(각 3~6글자). "
        "장식용 사진은 금지합니다. "
        "이미지는 단색/저채도/평면 도식 스타일로 간결하게 구성하세요. "
        "텍스트가 많은 카드는 이미지 없이 1열 텍스트 레이아웃으로 구성하고 "
        "폰트를 과도하게 축소하지 마세요."
    )


def _extract_pptx_text(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return ""

    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    tf = shape.text_frame
                    if tf and tf.text:
                        texts.append(tf.text)
            except Exception:
                continue
    return "\n".join(texts)


def _has_image_error_text(text: str) -> bool:
    if not text:
        return False
    low = text.lower()
    for phrase in _get_image_error_phrases():
        if phrase.lower() in low:
            return True
    return False


def _emu_to_pt(value: int) -> float:
    return value / 12700.0


def _build_overflow_report(pptx_path: str) -> Dict[str, Any]:
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"error": f"Failed to open PPTX: {e}"}

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    report_slides: List[Dict[str, Any]] = []
    overflow_shapes_total = 0
    slides_with_overflow = 0
    slides_with_small_font = 0
    min_font_global = None

    for idx, slide in enumerate(prs.slides, 1):
        overflow_shapes: List[Dict[str, Any]] = []
        min_font_pt = None

        for shape in slide.shapes:
            try:
                left = int(getattr(shape, "left", 0))
                top = int(getattr(shape, "top", 0))
                width = int(getattr(shape, "width", 0))
                height = int(getattr(shape, "height", 0))
                if (
                    left < 0
                    or top < 0
                    or left + width > slide_w
                    or top + height > slide_h
                ):
                    overflow_shapes.append(
                        {
                            "name": getattr(shape, "name", ""),
                            "type": str(getattr(shape, "shape_type", "")),
                            "leftPt": round(_emu_to_pt(left), 2),
                            "topPt": round(_emu_to_pt(top), 2),
                            "widthPt": round(_emu_to_pt(width), 2),
                            "heightPt": round(_emu_to_pt(height), 2),
                        }
                    )
            except Exception:
                continue

            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            size = run.font.size
                            if size is None:
                                continue
                            size_pt = float(size.pt)
                            if min_font_pt is None or size_pt < min_font_pt:
                                min_font_pt = size_pt
            except Exception:
                continue

        if overflow_shapes:
            slides_with_overflow += 1
            overflow_shapes_total += len(overflow_shapes)

        if min_font_pt is not None and min_font_pt < 10:
            slides_with_small_font += 1

        if min_font_pt is not None:
            if min_font_global is None or min_font_pt < min_font_global:
                min_font_global = min_font_pt

        report_slides.append(
            {
                "slideIndex": idx,
                "overflowShapes": overflow_shapes,
                "minFontPt": None if min_font_pt is None else round(min_font_pt, 2),
            }
        )

    summary = {
        "overflowSlides": slides_with_overflow,
        "overflowShapes": overflow_shapes_total,
        "slidesWithSmallFont": slides_with_small_font,
        "minFontPt": None if min_font_global is None else round(min_font_global, 2),
    }

    return {"summary": summary, "slides": report_slides}


def _parse_env_folder_ids() -> Optional[List[str]]:
    raw = os.environ.get("GAMMA_FOLDER_IDS")
    if not raw:
        return None
    items = [p.strip() for p in raw.split(",") if p.strip()]
    return items if items else None


def _resolve_image_style(image_source: str) -> Optional[str]:
    if image_source != "aiGenerated":
        return None
    env_style = os.environ.get("GAMMA_IMAGE_STYLE")
    if env_style:
        return env_style.strip()
    return DEFAULT_IMAGE_STYLE


def _resolve_image_model(image_source: str) -> Optional[str]:
    if image_source != "aiGenerated":
        return None
    env_model = os.environ.get("GAMMA_IMAGE_MODEL")
    if env_model:
        return env_model.strip()
    return None


def _run_generation(
    gamma_api_key: str,
    deck_text: str,
    wait: bool,
    num_cards_final: int,
    card_split: str,
    theme_id: Optional[str],
    folder_ids: Optional[List[str]],
    additional_instructions: Optional[str],
    image_source: str,
    image_style: Optional[str],
    image_model: Optional[str],
    mode: str,
    template_id: Optional[str],
    fallback: bool,
    template_strict: bool,
    timeout_sec: int,
    poll_interval_sec: int,
    warnings_in: Optional[List[str]] = None,
) -> Dict[str, Any]:
    warnings = list(warnings_in or [])
    mode_used = mode

    if mode == "template":
        if not template_id:
            if template_strict:
                raise RuntimeError("template_id is required for template_strict mode.")
            if fallback:
                warnings.append("Template mode requested but GAMMA_TEMPLATE_ID is not set. Falling back to generate.")
                mode_used = "generate"
            else:
                raise RuntimeError("template_id is required for template mode.")
        else:
            try:
                prompt_text = deck_text
                extra = _image_policy_instruction()
                if additional_instructions:
                    extra = f"{extra}\n{additional_instructions.strip()}"
                prompt_text = f"{prompt_text}\n\n[추가 지시]\n{extra}"
                payload_t = build_template_payload(
                    gamma_id=template_id,
                    prompt=prompt_text,
                    theme_id=theme_id,
                    folder_ids=folder_ids,
                    export_as="pptx",
                    image_style=image_style,
                    image_model=image_model,
                )
                generation_id = create_generation_from_template(gamma_api_key, payload_t)
            except Exception as e:
                if template_strict:
                    raise RuntimeError(f"Template mode failed (strict): {e}")
                if fallback:
                    warnings.append(f"Template mode failed: {e}. Falling back to generate.")
                    mode_used = "generate"
                else:
                    raise

    if mode_used != "template":
        extra_instructions = _image_policy_instruction()
        if additional_instructions:
            extra_instructions = f"{extra_instructions}\n{additional_instructions.strip()}"
        payload = build_gamma_payload(
            input_text=deck_text,
            card_split=card_split,
            num_cards=num_cards_final,
            theme_id=theme_id,
            folder_ids=folder_ids,
            additional_instructions=extra_instructions,
            image_source=image_source,
            image_style=image_style,
            image_model=image_model,
        )
        generation_id = create_generation(gamma_api_key, payload)

    if not wait:
        return {
            "status": "pending",
            "generationId": generation_id,
            "warnings": warnings,
            "modeUsed": mode_used,
        }

    data, status = poll_generation(
        gamma_api_key,
        generation_id,
        timeout_sec=timeout_sec,
        poll_interval_sec=poll_interval_sec,
    )

    if status == "timeout":
        return {
            "status": "pending",
            "generationId": generation_id,
            "warnings": warnings,
            "modeUsed": mode_used,
        }

    os.makedirs(RESP_DIR, exist_ok=True)
    resp_path = os.path.join(RESP_DIR, f"{generation_id}.json")
    with open(resp_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    gamma_url = data.get("gammaUrl")

    if str(data.get("status", "")).lower() != "completed":
        return {
            "status": data.get("status", "unknown"),
            "generationId": generation_id,
            "gammaUrl": gamma_url,
            "warnings": warnings,
            "modeUsed": mode_used,
        }

    pptx_url, url_warnings = choose_pptx_url(data)
    warnings.extend(url_warnings)
    if not pptx_url:
        return {
            "status": "completed",
            "generationId": generation_id,
            "gammaUrl": gamma_url,
            "warnings": warnings,
            "modeUsed": mode_used,
        }

    os.makedirs(EXPORT_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = os.path.join(EXPORT_DIR, f"gamma_{generation_id}_{timestamp}.pptx")
    download_pptx(pptx_url, out_path)

    return {
        "status": "completed",
        "generationId": generation_id,
        "gammaUrl": gamma_url,
        "pptxUrl": pptx_url,
        "pptxPath": out_path,
        "warnings": warnings,
        "modeUsed": mode_used,
    }


def generate_gamma_pptx(
    pdf_path: str,
    wait: bool = True,
    num_cards: Optional[int] = None,
    card_split: str = "inputTextBreaks",
    theme_id: Optional[str] = None,
    theme_id_from_request: bool = False,
    folder_ids: Optional[List[str]] = None,
    additional_instructions: Optional[str] = None,
    image_source: str = "aiGenerated",
    timeout_sec: int = 300,
    poll_interval_sec: int = 3,
    mode: str = "generate",
    template_id: Optional[str] = None,
    fallback: bool = True,
    template_strict: bool = False,
    condense_auto: bool = False,
    overflow_report: bool = False,
) -> Dict[str, Any]:
    gamma_api_key = os.environ.get("GAMMA_API_KEY")
    if not gamma_api_key:
        raise RuntimeError("GAMMA_API_KEY is not set.")

    raw_text = parse_pdf_to_text(pdf_path)
    if not raw_text.strip():
        raise RuntimeError("No text extracted from PDF.")

    target_chars = _compute_condense_target(len(raw_text), mode, condense_auto)
    rfp_text = _condense_if_needed(raw_text, target_chars)

    if mode == "template":
        if num_cards is not None and int(num_cards) != TEMPLATE_SLIDE_COUNT:
            warnings = [f"Template mode uses fixed {TEMPLATE_SLIDE_COUNT} slides; num_cards ignored."]
        else:
            warnings = []
        num_cards_final = TEMPLATE_SLIDE_COUNT
        try:
            template_json = _generate_template_json(rfp_text)
            slides = _parse_template_slides(template_json)
            deck_text = _slides_to_deck_text(slides)
        except Exception:
            deck_text = _generate_deck_text(rfp_text, num_cards_final)
    else:
        if num_cards is None:
            recommended = _recommend_card_count(rfp_text, MIN_CARDS, MAX_CARDS)
            num_cards_final = _clamp_cards(recommended, MIN_CARDS, MAX_CARDS)
        else:
            num_cards_final = int(num_cards)
        warnings = []
        deck_text = _generate_deck_text(rfp_text, num_cards_final)

    deck_text = _clean_deck_text(deck_text)

    if mode != "template":
        if theme_id is None:
            theme_id = os.environ.get("GAMMA_THEME_ID") or None
    else:
        if not theme_id_from_request:
            theme_id = None

    if folder_ids is None:
        folder_ids = _parse_env_folder_ids()

    warnings = warnings if isinstance(warnings, list) else []
    image_style = _resolve_image_style(image_source)
    image_model = _resolve_image_model(image_source)

    result = _run_generation(
        gamma_api_key=gamma_api_key,
        deck_text=deck_text,
        wait=wait,
        num_cards_final=num_cards_final,
        card_split=card_split,
        theme_id=theme_id,
        folder_ids=folder_ids,
        additional_instructions=additional_instructions,
        image_source=image_source,
        image_style=image_style,
        image_model=image_model,
        mode=mode,
        template_id=template_id,
        fallback=fallback,
        template_strict=template_strict,
        timeout_sec=timeout_sec,
        poll_interval_sec=poll_interval_sec,
        warnings_in=warnings,
    )

    if (
        wait
        and result.get("status") == "completed"
        and result.get("pptxPath")
        and image_source == "aiGenerated"
    ):
        pptx_text = _extract_pptx_text(result["pptxPath"])
        if _has_image_error_text(pptx_text):
            warn_list = list(result.get("warnings") or [])
            warn_list.append("Image error detected; fallback to placeholder.")

            fallback_mode = result.get("modeUsed") or mode
            if fallback_mode == "template":
                if template_strict:
                    warn_list.append("Template strict enabled; image error fallback skipped.")
                    result["warnings"] = warn_list
                    return result
                warn_list.append("Template mode image error; fallback to generate + placeholder.")
                fallback_mode = "generate"

            result = _run_generation(
                gamma_api_key=gamma_api_key,
                deck_text=deck_text,
                wait=wait,
                num_cards_final=num_cards_final,
                card_split=card_split,
                theme_id=theme_id,
                folder_ids=folder_ids,
                additional_instructions=additional_instructions,
                image_source="placeholder",
                image_style=None,
                image_model=None,
                mode=fallback_mode,
                template_id=template_id,
                fallback=fallback,
                template_strict=template_strict,
                timeout_sec=timeout_sec,
                poll_interval_sec=poll_interval_sec,
                warnings_in=warn_list,
            )

    if wait and result.get("status") == "completed" and result.get("pptxPath") and overflow_report:
        report = _build_overflow_report(result["pptxPath"])
        os.makedirs(REPORT_DIR, exist_ok=True)
        report_path = os.path.join(REPORT_DIR, f"{result.get('generationId')}.json")
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        result["overflowReport"] = report
        result["overflowReportPath"] = report_path

        summary = report.get("summary") if isinstance(report, dict) else {}
        if isinstance(summary, dict):
            overflow_shapes = summary.get("overflowShapes") or 0
            small_fonts = summary.get("slidesWithSmallFont") or 0
            if overflow_shapes or small_fonts:
                warn_list = list(result.get("warnings") or [])
                warn_list.append(
                    f"Overflow report: shapes={overflow_shapes}, smallFontSlides={small_fonts}."
                )
                result["warnings"] = warn_list

    return result
