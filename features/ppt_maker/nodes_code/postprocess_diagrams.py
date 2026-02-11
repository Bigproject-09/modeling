from __future__ import annotations

import re
from typing import Any, Dict, List, Optional, Set

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


# ✅ 사용자 요구 순서(고정)
SECTION_ORDER = [
    "기관 소개",
    "사업 개요",
    "연구 목표",
    "연구 필요성",
    "연구 내용",
    "추진 계획",
    "기대 효과",
    "활용 계획",
]


def _norm(s: Any) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


def _remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _delete_slide(prs: Presentation, slide_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    sldId = sldIdLst[slide_index]
    rId = sldId.rId
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _clear_slide(slide) -> None:
    for sh in list(slide.shapes):
        _remove_shape(sh)


def _extract_best_title(state: Dict[str, Any]) -> str:
    deck = (state.get("deck_json") or {})
    slides = deck.get("slides") or []
    if slides and isinstance(slides[0], dict):
        first_title = _norm(slides[0].get("slide_title") or "")
        if first_title:
            return first_title
    title = _norm(deck.get("deck_title") or "")
    if title and title != "(과제명 미기재)":
        return title
    # fallback
    return title or "(과제명 미기재)"


def _write_cover(slide, title: str) -> None:
    _clear_slide(slide)

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.0))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)


def _write_agenda(slide) -> None:
    _clear_slide(slide)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_t = title_box.text_frame
    tf_t.clear()
    p0 = tf_t.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "목차"
    r0.font.size = Pt(32)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(20, 20, 20)

    body = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(5.2))
    tf = body.text_frame
    tf.clear()

    for i, item in enumerate(SECTION_ORDER, 1):
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = f"{i}. {item}"
        p.level = 0
        if p.runs:
            p.runs[0].font.size = Pt(22)
        else:
            r = p.add_run()
            r.font.size = Pt(22)


def _slides_with_structured_visuals(deck_json: Dict[str, Any]) -> Set[int]:
    """
    deck_json의 슬라이드 중 TABLE/DIAGRAM/CHART 스펙이 있는 슬라이드 index(0-based)를 반환.
    이런 슬라이드는 Gamma가 그림(PICTURE)로 렌더링할 수 있으니 PICTURE를 보존해야 함.
    """
    keep: Set[int] = set()
    slides = (deck_json or {}).get("slides") or []
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue
        table_md = _norm(s.get("TABLE_MD") or "")
        diagram = _norm(s.get("DIAGRAM_SPEC_KO") or "")
        chart = _norm(s.get("CHART_SPEC_KO") or "")
        if table_md or diagram or chart:
            keep.add(i)
    return keep


def _slide_has_structured_spec(slide_spec: Dict[str, Any]) -> bool:
    if not isinstance(slide_spec, dict):
        return False
    return bool(
        _norm(slide_spec.get("TABLE_MD") or "")
        or _norm(slide_spec.get("DIAGRAM_SPEC_KO") or "")
        or _norm(slide_spec.get("CHART_SPEC_KO") or "")
    )


def _remove_visual_placeholders(prs: Presentation, keep_picture_slide_idxs: Set[int]) -> int:
    """
    - AI 이미지/placeholder만 제거
    - 단, 구조화된 시각자료(표/차트/다이어그램) 스펙이 있는 슬라이드는 PICTURE 보존
    """
    removed = 0
    for si, slide in enumerate(prs.slides):
        for sh in list(slide.shapes):
            # 1) 실제 그림
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # ✅ 표/도표/다이어그램이 그림으로 들어온 슬라이드는 보존
                if si in keep_picture_slide_idxs:
                    continue
                _remove_shape(sh)
                removed += 1
                continue

            # 2) placeholder (picture)
            if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                ph = getattr(sh, "placeholder_format", None)
                ph_type = str(getattr(ph, "type", "")).lower() if ph else ""
                nm = (getattr(sh, "name", "") or "").lower()
                if ("picture" in ph_type) or ("pic" in ph_type) or ("image" in nm) or ("picture" in nm) or ("이미지" in nm):
                    _remove_shape(sh)
                    removed += 1
                    continue

            # 3) 도형으로 만들어진 '이미지 자리'만 제거 (실제 다이어그램 도형은 건드리지 않도록 이름 기반으로만)
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                nm = (getattr(sh, "name", "") or "").lower()
                if any(k in nm for k in ["picture placeholder", "image placeholder", "picture", "image"]):
                    _remove_shape(sh)
                    removed += 1
                    continue

    return removed


def _trim_ending_slides(prs: Presentation) -> None:
    """
    마지막에 감사/Q&A가 여러 장 생기거나, 쓸데없는 엔딩이 붙는 경우 정리.
    - '감사합니다'는 1장만 유지
    - '추가 정보/문의/연락처/회사 소개' 류 제거
    """
    bad_keywords = ["추가 정보", "문의", "연락처", "회사 소개", "contact", "thank you", "thanks"]
    thanks_idx: List[int] = []

    for i, slide in enumerate(prs.slides):
        txt = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt += " " + (sh.text_frame.text or "")
        t = _norm(txt).lower()
        if "감사합니다" in t:
            thanks_idx.append(i)

    # 감사합니다가 2장 이상이면 마지막 1장만 남김
    if len(thanks_idx) > 1:
        for idx in reversed(thanks_idx[:-1]):
            _delete_slide(prs, idx)

    # 마지막 쪽에서 bad 키워드 슬라이드 제거 (보수적으로: 마지막 5장 범위)
    n = len(prs.slides)
    for idx in reversed(range(max(0, n - 5), n)):
        slide = prs.slides[idx]
        txt = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt += " " + (sh.text_frame.text or "")
        t = _norm(txt)
        if any(k.lower() in t.lower() for k in bad_keywords):
            # 단, '감사합니다'는 남김
            if "감사합니다" in t:
                continue
            _delete_slide(prs, idx)


def postprocess_diagrams(pptx_path: str, deck_json: Dict[str, Any], state: Optional[Dict[str, Any]] = None) -> str:
    """
    Gamma 결과 PPTX 후처리
    - 표지/목차 강제 재작성 (1~2번 슬라이드)
    - AI 이미지/placeholder 제거 (단, 표/도표/다이어그램은 보존)
    - 엔딩 슬라이드 정리
    """
    prs = Presentation(pptx_path)

    # ✅ deck_json 기반으로 "그림 보존해야 하는 슬라이드" 계산
    keep_picture_slide_idxs = _slides_with_structured_visuals(deck_json)
    deck_slides = (deck_json or {}).get("slides") or []

    # 1) 표지/목차 재작성 (Gamma가 이상하게 만들어도 여기서 고정)
    if len(prs.slides) >= 1:
        cover_spec = deck_slides[0] if len(deck_slides) >= 1 and isinstance(deck_slides[0], dict) else {}
        if not _slide_has_structured_spec(cover_spec):
            title = _extract_best_title(state or {"deck_json": deck_json})
            _write_cover(prs.slides[0], title)

    if len(prs.slides) >= 2:
        agenda_spec = deck_slides[1] if len(deck_slides) >= 2 and isinstance(deck_slides[1], dict) else {}
        if not _slide_has_structured_spec(agenda_spec):
            _write_agenda(prs.slides[1])

    # 2) AI 이미지/placeholder 제거(표/도표/다이어그램은 보존)
    _remove_visual_placeholders(prs, keep_picture_slide_idxs)

    # 3) 엔딩 정리
    _trim_ending_slides(prs)

    prs.save(pptx_path)
    return pptx_path


def postprocess_diagrams_node(state: Dict[str, Any]) -> Dict[str, Any]:
    pptx_path = state.get("final_ppt_path") or state.get("gamma_ppt_path")
    deck_json = state.get("deck_json") or {}
    if not pptx_path:
        return state
    postprocess_diagrams(pptx_path, deck_json, state=state)
    return state
