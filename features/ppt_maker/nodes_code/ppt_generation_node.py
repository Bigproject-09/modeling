# import os
# import json
# from typing import List, Dict
# from datetime import datetime
# from dotenv import load_dotenv
# from google import genai
# from google.genai import types
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from io import BytesIO
# from PIL import Image

# load_dotenv()

# # 색상 테마
# COLOR_THEMES = {
#     "ocean_polar": {
#         "primary": RGBColor(6, 90, 130),
#         "secondary": RGBColor(28, 114, 147),
#         "accent": RGBColor(33, 41, 92),
#         "text_main": RGBColor(30, 30, 30),
#         "text_light": RGBColor(100, 100, 100),
#         "text_inv": RGBColor(255, 255, 255),
#         "bg_light": RGBColor(248, 249, 251),
#         "card_bg": RGBColor(240, 245, 250)
#     }
# }

# CURRENT_COLORS = COLOR_THEMES["ocean_polar"]
# SLIDE_WIDTH, SLIDE_HEIGHT = 10, 7.5
# MAX_Y = 7.0

# def select_color_theme(title: str) -> str:
#     title_lower = title.lower()
#     if any(w in title_lower for w in ['해양','극지','기후']):
#         return "ocean_polar"
#     return "ocean_polar"

# # LLM 디자인 가이드
# def get_ppt_design_guide(project_title: str, slides_data: List[dict]) -> Dict:
#     try:
#         api_key = os.environ.get("GEMINI_API_KEY")
#         if not api_key:
#             return get_fallback_design(slides_data)
            
#         client = genai.Client(api_key=api_key)
        
#         slides_summary = []
#         for idx, slide in enumerate(slides_data, 1):
#             items = slide.get("items", [])
#             slides_summary.append({
#                 "index": idx,
#                 "title": slide.get("title", "")[:40],
#                 "item_count": len(items)
#             })
        
#         prompt = f"""전문 디자이너로서 PPT 디자인을 결정하세요.

# **프로젝트:** {project_title}
# **슬라이드:** {len(slides_data)}개
# **구성:** {json.dumps(slides_summary, ensure_ascii=False)}

# **결정사항:**
# 1. template.header_style: "minimal" (얇은 라인) | "bold" (배경 헤더)
# 2. template.spacing: "normal" | "spacious"
# 3. 각 슬라이드의 layout: "vertical_list" | "horizontal_cards" | "highlight_boxes"
# 4. 각 항목의 위치 (x, y, width, height in inches)

# **제약조건:**
# - 제목이 길면 2줄 가능 (y=0.3~1.1 영역)
# - 내용 시작: y >= 1.4
# - 최대: y + height <= 7.0
# - 박스 내부 여백: 0.3인치
# - 항목 간 간격: 최소 0.2인치

# JSON:
# {{
#   "template": {{"header_style": "minimal", "spacing": "normal"}},
#   "slides": [
#     {{
#       "index": 1,
#       "layout": "vertical_list",
#       "items": [
#         {{"x": 1.0, "y": 1.5, "width": 8.0, "height": 1.0}},
#         {{"x": 1.0, "y": 2.7, "width": 8.0, "height": 1.0}}
#       ]
#     }}
#   ]
# }}"""
        
#         response = client.models.generate_content(
#             model="gemini-2.5-flash",
#             contents=prompt,
#             config=types.GenerateContentConfig(response_mime_type="application/json", temperature=0.3)
#         )
        
#         result = json.loads(response.text)
#         result = validate_design(result, slides_data)
        
#         print(f"\n  [LLM 디자인]")
#         print(f"    템플릿: {result['template']['header_style']} | {result['template']['spacing']}")
#         print(f"    슬라이드: {len(result['slides'])}개 설계")
        
#         return result
        
#     except Exception as e:
#         print(f"\n  [LLM 실패] {e}")
#         return get_fallback_design(slides_data)

# def validate_design(design: Dict, slides_data: List[dict]) -> Dict:
#     """디자인 유효성 검증"""
#     if len(design.get('slides', [])) != len(slides_data):
#         design['slides'] = []
#         for idx, slide in enumerate(slides_data, 1):
#             items = slide.get('items', [])
#             design['slides'].append({
#                 "index": idx,
#                 "layout": "vertical_list",
#                 "items": [{"x": 1.0, "y": 1.5 + i*1.1, "width": 8.0, "height": 0.9} for i in range(min(len(items), 5))]
#             })
    
#     for slide_design in design['slides']:
#         for item in slide_design.get('items', []):
#             if item['y'] + item['height'] > MAX_Y:
#                 item['height'] = MAX_Y - item['y'] - 0.1
#             item['height'] = max(0.4, min(item['height'], 2.0))
#             item['width'] = max(3.0, min(item['width'], 9.0))
    
#     return design

# def get_fallback_design(slides_data: List[dict]) -> Dict:
#     slides = []
#     for idx, slide in enumerate(slides_data, 1):
#         items = slide.get('items', [])
#         slides.append({
#             "index": idx,
#             "layout": "vertical_list",
#             "items": [{"x": 1.0, "y": 1.5 + i*1.1, "width": 8.0, "height": 0.9} for i in range(min(len(items), 5))]
#         })
    
#     return {
#         "template": {"header_style": "minimal", "spacing": "normal"},
#         "slides": slides
#     }

# # 표지
# def render_cover(slide, title: str, researcher: str):
#     """간단하고 깔끔한 표지"""
#     # 배경
#     bg = slide.shapes.add_shape(1, 0, 0, Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
#     bg.fill.solid()
#     bg.fill.fore_color.rgb = CURRENT_COLORS["primary"]
#     bg.line.fill.background()
#     bg.shadow.inherit = False
    
#     # 좌측 바
#     bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), Inches(SLIDE_HEIGHT))
#     bar.fill.solid()
#     bar.fill.fore_color.rgb = CURRENT_COLORS["accent"]
#     bar.line.fill.background()
#     bar.shadow.inherit = False
    
#     # 제목
#     title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(8.0), Inches(2.2))
#     tf = title_box.text_frame
#     tf.word_wrap = True
#     tf.vertical_anchor = 1
#     p = tf.paragraphs[0]
#     p.text = title
#     p.font.size = Pt(36)
#     p.font.bold = True
#     p.font.color.rgb = CURRENT_COLORS["text_inv"]
#     p.alignment = PP_ALIGN.LEFT
#     p.line_spacing = 1.25
    
#     # 연구자
#     researcher_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(8.0), Inches(0.4))
#     rf = researcher_box.text_frame
#     rp = rf.paragraphs[0]
#     rp.text = f"{researcher}  |  연구책임자"
#     rp.font.size = Pt(15)
#     rp.font.color.rgb = CURRENT_COLORS["text_inv"]
#     rp.alignment = PP_ALIGN.LEFT

# # 헤더
# def render_header(slide, title: str, header_style: str):
#     """제목 길이 자동 조정"""
#     title_length = len(title)
#     font_size = 20 if title_length > 45 else (23 if title_length > 30 else 26)
    
#     if header_style == "bold":
#         header_bg = slide.shapes.add_shape(1, 0, 0, Inches(SLIDE_WIDTH), Inches(1.0))
#         header_bg.fill.solid()
#         header_bg.fill.fore_color.rgb = CURRENT_COLORS["bg_light"]
#         header_bg.line.fill.background()
#         header_bg.shadow.inherit = False
        
#         title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.65))
#     else:
#         title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.75))
    
#     tf = title_box.text_frame
#     tf.word_wrap = True
#     p = tf.paragraphs[0]
#     p.text = title
#     p.font.size = Pt(font_size)
#     p.font.bold = True
#     p.font.color.rgb = CURRENT_COLORS["primary"]
#     p.line_spacing = 1.15
    
#     if header_style != "bold":
#         line = slide.shapes.add_shape(1, Inches(0.4), Inches(1.08), Inches(9.2), Inches(0.015))
#         line.fill.solid()
#         line.fill.fore_color.rgb = CURRENT_COLORS["primary"]
#         line.line.fill.background()
#         line.shadow.inherit = False

# # 콘텐츠 렌더링
# def render_content(slide, items: List[dict], layout_design: Dict, layout_type: str):
#     """LLM 배치대로 렌더링 + 이미지 지원 (position: left/right/top/bottom)"""
#     items_positions = layout_design.get('items', [])
    
#     for idx, (item, pos) in enumerate(zip(items, items_positions)):
#         if idx >= len(items_positions):
#             break
            
#         subtitle = item.get("subtitle", "")
#         content = item.get("content", "")
#         image_path = item.get("image_path", None)  # 앞 노드에서 전달
#         image_position = item.get("image_position", "right")  # left/right/top/bottom
        
#         x = pos.get('x', 1.0)
#         y = pos.get('y', 1.5)
#         width = pos.get('width', 8.0)
#         height = pos.get('height', 1.0)
        
#         if y + height > MAX_Y:
#             height = MAX_Y - y - 0.1
#             if height < 0.3:
#                 break
        
#         if layout_type == "highlight_boxes":
#             # 강조 박스
#             color_bar = slide.shapes.add_shape(1, Inches(x - 0.15), Inches(y), Inches(0.06), Inches(height))
#             color_bar.fill.solid()
#             color_bar.fill.fore_color.rgb = CURRENT_COLORS["primary"]
#             color_bar.line.fill.background()
#             color_bar.shadow.inherit = False
            
#             bg_box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
#             bg_box.fill.solid()
#             bg_box.fill.fore_color.rgb = CURRENT_COLORS["card_bg"]
#             bg_box.line.fill.background()
#             bg_box.shadow.inherit = False
            
#             # 텍스트 영역 (내부 여백 확보)
#             text_x = x + 0.25
#             text_y = y + 0.15
#             text_w = width - 0.5
#             text_h = height - 0.3
            
#             # 이미지가 있으면 위치에 따라 배치
#             if image_path and os.path.exists(image_path):
#                 img_gap = 0.2  # 이미지-텍스트 간격
                
#                 if image_position == "right":
#                     # 이미지 우측
#                     img_w = text_w * 0.35
#                     img_h = text_h
#                     slide.shapes.add_picture(
#                         image_path, 
#                         Inches(text_x + text_w - img_w), 
#                         Inches(text_y), 
#                         width=Inches(img_w), 
#                         height=Inches(img_h)
#                     )
#                     text_w = text_w * 0.60  # 텍스트 영역 축소
                
#                 elif image_position == "left":
#                     # 이미지 좌측
#                     img_w = text_w * 0.35
#                     img_h = text_h
#                     slide.shapes.add_picture(
#                         image_path, 
#                         Inches(text_x), 
#                         Inches(text_y), 
#                         width=Inches(img_w), 
#                         height=Inches(img_h)
#                     )
#                     text_x = text_x + img_w + img_gap
#                     text_w = text_w * 0.60
                
#                 elif image_position == "top":
#                     # 이미지 상단
#                     img_h = text_h * 0.4
#                     img_w = text_w
#                     slide.shapes.add_picture(
#                         image_path, 
#                         Inches(text_x), 
#                         Inches(text_y), 
#                         width=Inches(img_w), 
#                         height=Inches(img_h)
#                     )
#                     text_y = text_y + img_h + img_gap
#                     text_h = text_h * 0.55
                
#                 elif image_position == "bottom":
#                     # 이미지 하단
#                     img_h = text_h * 0.4
#                     img_w = text_w
#                     slide.shapes.add_picture(
#                         image_path, 
#                         Inches(text_x), 
#                         Inches(text_y + text_h - img_h), 
#                         width=Inches(img_w), 
#                         height=Inches(img_h)
#                     )
#                     text_h = text_h * 0.55
            
#             text_box = slide.shapes.add_textbox(Inches(text_x), Inches(text_y), Inches(text_w), Inches(text_h))
#             tf = text_box.text_frame
#             tf.word_wrap = True
            
#             if subtitle:
#                 p = tf.paragraphs[0]
#                 p.text = subtitle
#                 p.font.size = Pt(18)
#                 p.font.bold = True
#                 p.font.color.rgb = CURRENT_COLORS["primary"]
#                 p.space_after = Pt(4)
            
#             if content:
#                 cp = tf.add_paragraph()
#                 text = content.split('\n')[0].strip().lstrip('•-·')[:100]
#                 cp.text = text
#                 cp.font.size = Pt(14)
#                 cp.font.color.rgb = CURRENT_COLORS["text_light"]
        
#         else:  # vertical_list
#             # 리스트
#             if subtitle:
#                 dot = slide.shapes.add_shape(2, Inches(x - 0.25), Inches(y + 0.05), Inches(0.1), Inches(0.1))
#                 dot.fill.solid()
#                 dot.fill.fore_color.rgb = CURRENT_COLORS["primary"]
#                 dot.line.fill.background()
#                 dot.shadow.inherit = False
                
#                 st = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.3))
#                 st.text_frame.word_wrap = True
#                 st.text_frame.paragraphs[0].text = subtitle
#                 st.text_frame.paragraphs[0].font.size = Pt(17)
#                 st.text_frame.paragraphs[0].font.bold = True
#                 st.text_frame.paragraphs[0].font.color.rgb = CURRENT_COLORS["text_main"]
            
#             if content:
#                 ct_y = y + 0.35 if subtitle else y
#                 ct_h = height - 0.35 if subtitle else height
                
#                 ct = slide.shapes.add_textbox(Inches(x), Inches(ct_y), Inches(width), Inches(ct_h))
#                 ct.text_frame.word_wrap = True
                
#                 lines = [l.strip().lstrip('•-·') for l in content.split('\n') if l.strip()]
#                 for line in lines[:2]:
#                     p = ct.text_frame.add_paragraph() if ct.text_frame.text else ct.text_frame.paragraphs[0]
#                     p.text = line[:90]
#                     p.font.size = Pt(14)
#                     p.font.color.rgb = CURRENT_COLORS["text_light"]

# # 메인
# def generate_ppt_node(state: Dict, output_dir: str, project_title: str = None, researcher_name: str = "박OO 교수") -> dict:
#     try:
#         prs = Presentation()
#         prs.slide_width = Inches(SLIDE_WIDTH)
#         prs.slide_height = Inches(SLIDE_HEIGHT)
        
#         slides_data = state.get("slides", [])
        
#         if project_title is None:
#             project_title = state.get("analyzed_json", {}).get("project_summary", {}).get("title", "연구 제안서")
        
#         print(f"\n[PPT 생성]")
#         print(f"  제목: {project_title}")
#         print(f"  슬라이드: {len(slides_data)}개")
        
#         theme_key = select_color_theme(project_title)
#         global CURRENT_COLORS
#         CURRENT_COLORS = COLOR_THEMES[theme_key]
        
#         design_guide = get_ppt_design_guide(project_title, slides_data)
        
#         # 표지
#         cover = prs.slides.add_slide(prs.slide_layouts[6])
#         render_cover(cover, project_title, researcher_name)
#         print(f"\n  ✓ 표지")
        
#         # 내용
#         template = design_guide['template']
#         slides_designs = design_guide['slides']
        
#         for idx, s_data in enumerate(slides_data, 1):
#             slide = prs.slides.add_slide(prs.slide_layouts[6])
#             title = s_data.get("title", "")
#             items = s_data.get("items", [])
            
#             print(f"  [{idx}] {title[:30]}...")
            
#             render_header(slide, title, template['header_style'])
            
#             layout_design = slides_designs[idx - 1] if idx <= len(slides_designs) else slides_designs[0]
#             layout_type = layout_design.get('layout', 'vertical_list')
            
#             render_content(slide, items, layout_design, layout_type)
        
#         os.makedirs(output_dir, exist_ok=True)
#         filename = f"PPT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
#         output_path = os.path.join(output_dir, filename)
#         prs.save(output_path)
        
#         print(f"\n[완료] {output_path}\n")
#         return {"ppt_path": output_path}
        
#     except Exception as e:
#         print(f"\n[실패] {e}\n")
#         import traceback
#         traceback.print_exc()
#         return {"ppt_path": ""}

# # 테스트
# def test_generate_ppt_node():
#     """이미지 포함 테스트 - state에서 image_path와 position 받아옴"""
    
#     # ========================================
#     # 테스트용 이미지 저장 위치 안내
#     # ========================================
#     # 실제 사용 시: 이미지 생성 노드에서 이미지를 생성하고
#     # 경로를 C:\big_project\modeling\data\images\ 에 저장
#     # 예: C:\big_project\modeling\data\images\slide1_image.png
    
#     # 테스트용으로는 업로드된 이미지 사용
#     test_image_path = "./data/images/test-image.png"
    
#     # 실제 구조 시뮬레이션
#     test_slides = [
#         {
#             "page_number": 1,
#             "title": "연구개발 목표 및 핵심 기술",
#             "items": [
#                 {
#                     "subtitle": "연구 목표", 
#                     "content": "세계 3위권 예측 정확도 확보를 통한 국제 경쟁력 강화\n계절~10년 규모 장기 예측 기술 확보"
#                 },
#                 {
#                     "subtitle": "초고해상도 모델링", 
#                     "content": "해상도 혁신: 기존 1도(100km) → 1/10도(10km)로 10배 향상\n중규모 에디(Mesoscale Eddy) 명시적 계산 가능",
#                     "image_path": test_image_path,  # 앞 노드에서 전달
#                     "image_position": "right"       # 앞 노드에서 전달 (left/right/top/bottom)
#                 },
#                 {
#                     "subtitle": "AI 기반 오차 보정", 
#                     "content": "딥러닝 기반 편향 보정(Bias Correction) 시스템\n모델 체계적 오차를 AI가 학습하여 실시간 보정",
#                     "image_path": test_image_path,  
#                     "image_position": "left"        # 다른 위치 테스트
#                 },
#                 {
#                     "subtitle": "통합 결합 모델", 
#                     "content": "해양(MOM6) + 빙권(CICE6) + 생태계(NPZD) 완전 결합\n물리-생태계 양방향 피드백 구현"
#                 }
#             ]
#         }
#     ]
    
#     test_state = {
#         "slides": test_slides,
#         "analyzed_json": {
#             "project_summary": {
#                 "title": "차세대 해양·극지 환경 기후예측시스템 개발"
#             }
#         }
#     }
    
#     result = generate_ppt_node(
#         test_state, 
#         r"C:\big_project\modeling\data\pptx", 
#         "차세대 해양·극지 환경 및 생태계 기후예측시스템 개발",
#         "박철수 교수"
#     )
    
#     print(f"\n{'='*60}")
#     if result['ppt_path']:
#         print(f"테스트 성공!")
#         print(f"파일: {result['ppt_path']}")
#         print(f"\n이미지 배치 테스트:")
#         print(f"  - 항목 2: '초고해상도 모델링' (이미지 우측)")
#         print(f"  - 항목 3: 'AI 기반 오차 보정' (이미지 좌측)")
#         print(f"\n실제 사용 시:")
#         print(f"  이미지 저장 경로: C:\\big_project\\modeling\\data\\images\\")
#         print(f"  예시: slide1_image.png, slide2_diagram.png 등")
#     else:
#         print(f"테스트 실패")
#     print(f"{'='*60}\n")

# if __name__ == "__main__":
#     test_generate_ppt_node()
import os
import json
import re
from typing import List, Dict, Optional
from datetime import datetime
from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import base64

load_dotenv()

# ============================================
# 색상 테마 정의 (간소화)
# ============================================

COLOR_THEMES = {
    "ocean_marine": {
        "keywords": ["해양", "바다", "수산", "극지", "빙하", "해류"],
        "primary": "#065A82", "secondary": "#1C7293", "accent": "#21295C",
        "text_main": "#1E1E1E", "text_light": "#646464", "text_inverse": "#FFFFFF",
        "bg_light": "#F8F9FB", "card_bg": "#F0F5FA"
    },
    "biotech_health": {
        "keywords": ["생명", "의료", "바이오", "건강", "질병", "치료", "유전자"],
        "primary": "#0C6A4E", "secondary": "#10A37F", "accent": "#4A5568",
        "text_main": "#2D3748", "text_light": "#718096", "text_inverse": "#FFFFFF",
        "bg_light": "#F7FAFC", "card_bg": "#EDF7F4"
    },
    "energy_environment": {
        "keywords": ["에너지", "환경", "신재생", "태양광", "풍력", "기후", "탄소"],
        "primary": "#047857", "secondary": "#10B981", "accent": "#F59E0B",
        "text_main": "#1F2937", "text_light": "#6B7280", "text_inverse": "#FFFFFF",
        "bg_light": "#F0FDF4", "card_bg": "#ECFDF5"
    },
    "ai_tech": {
        "keywords": ["인공지능", "AI", "딥러닝", "머신러닝", "데이터", "알고리즘"],
        "primary": "#5B21B6", "secondary": "#7C3AED", "accent": "#EC4899",
        "text_main": "#1F2937", "text_light": "#6B7280", "text_inverse": "#FFFFFF",
        "bg_light": "#FAF5FF", "card_bg": "#F3E8FF"
    },
    "space_aerospace": {
        "keywords": ["우주", "항공", "위성", "로켓", "항공우주"],
        "primary": "#1E3A8A", "secondary": "#3B82F6", "accent": "#F97316",
        "text_main": "#111827", "text_light": "#6B7280", "text_inverse": "#FFFFFF",
        "bg_light": "#EFF6FF", "card_bg": "#DBEAFE"
    },
    "manufacturing": {
        "keywords": ["제조", "생산", "공정", "자동화", "스마트팩토리", "산업"],
        "primary": "#92400E", "secondary": "#D97706", "accent": "#475569",
        "text_main": "#1C1917", "text_light": "#78716C", "text_inverse": "#FFFFFF",
        "bg_light": "#FFFBEB", "card_bg": "#FEF3C7"
    },
    "finance": {
        "keywords": ["금융", "경제", "투자", "자산", "재무", "증권"],
        "primary": "#1E40AF", "secondary": "#3B82F6", "accent": "#059669",
        "text_main": "#1E293B", "text_light": "#475569", "text_inverse": "#FFFFFF",
        "bg_light": "#F0F9FF", "card_bg": "#E0F2FE"
    },
    "professional": {
        "keywords": ["연구", "개발", "프로젝트", "기술", "혁신"],  # 기본 테마
        "primary": "#1E40AF", "secondary": "#3B82F6", "accent": "#64748B",
        "text_main": "#1E293B", "text_light": "#64748B", "text_inverse": "#FFFFFF",
        "bg_light": "#F8FAFC", "card_bg": "#F1F5F9"
    }
}


def select_theme(state: Dict, project_title: str, project_description: str = "") -> str:
    """
    개선된 테마 선택 - state의 슬라이드 내용도 분석
    
    Args:
        state: 슬라이드 데이터 포함
        project_title: 프로젝트 제목
        project_description: 프로젝트 설명
    
    Returns:
        테마 키 (예: "ocean_marine")
    """
    # 1. 제목/설명
    text_parts = [project_title, project_description]
    
    # 2. state에서 추가 텍스트 추출
    slides_data = state.get("slides", [])
    
    # 슬라이드 제목들 (최대 3개)
    for slide in slides_data[:3]:
        text_parts.append(slide.get("title", ""))
    
    # 항목 부제목/내용 (처음 2개 슬라이드의 각 2개 항목)
    for slide in slides_data[:2]:
        for item in slide.get("items", [])[:2]:
            text_parts.append(item.get("subtitle", ""))
            text_parts.append(item.get("content", "")[:100])  # 내용은 100자만
    
    # 3. 전체 텍스트 합치기
    full_text = " ".join(text_parts).lower()
    
    # 4. 키워드 매칭
    scores = {}
    matched_keywords = {}
    
    for theme_key, theme in COLOR_THEMES.items():
        score = 0
        matched = []
        
        for keyword in theme["keywords"]:
            if keyword in full_text:
                score += 1
                matched.append(keyword)
        
        scores[theme_key] = score
        matched_keywords[theme_key] = matched
    
    # 5. 최고 점수 테마 선택
    best_theme = max(scores.items(), key=lambda x: x[1])
    theme_key = best_theme[0]
    theme_score = best_theme[1]
    
    # 6. 매칭 결과 출력
    if theme_score == 0:
        print(f"  [테마 선택] professional (기본) - 키워드 매칭 없음")
        return "professional"
    else:
        matched = matched_keywords[theme_key]
        print(f"  [테마 선택] {theme_key}")
        print(f"    매칭 키워드: {', '.join(matched[:5])} (총 {theme_score}개)")
    
    return theme_key


# ============================================
# LLM이 PPTX 구조 JSON 직접 생성
# ============================================

def generate_ppt_structure(slides_data: List[dict], project_title: str, theme: Dict) -> Dict:
    """LLM이 슬라이드 내용을 분석하여 최적의 PPTX 구조를 JSON으로 생성"""
    try:
        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            return generate_fallback_structure(slides_data, project_title, theme)

        client = genai.Client(api_key=api_key)

        # 내용 분석
        content_analysis = []
        for idx, slide in enumerate(slides_data, 1):
            items = slide.get("items", [])
            total_text = sum(len(item.get('subtitle', '') + item.get('content', '')) for item in items)
            content_analysis.append({
                "slide_num": idx,
                "title": slide.get("title", "")[:60],
                "item_count": len(items),
                "total_chars": total_text,
                "needs_space": total_text > 400
            })

        prompt = f"""당신은 전문 PPT 디자이너입니다. 슬라이드 내용을 분석하여 **최적의 레이아웃 구조를 JSON으로 생성**하세요.

**프로젝트:** {project_title}

**슬라이드 분석:**
{json.dumps(content_analysis, ensure_ascii=False, indent=2)}

**전체 데이터:**
{json.dumps(slides_data, ensure_ascii=False, indent=2)}

**색상:**
Primary: {theme['primary']}, Secondary: {theme['secondary']}, Accent: {theme['accent']}
Text: {theme['text_main']}, Light: {theme['text_light']}, Inverse: {theme['text_inverse']}
BG: {theme['bg_light']}, Card: {theme['card_bg']}

**슬라이드 크기:** 10인치 × 7.5인치
**단위:** 모든 위치/크기는 인치 (inch)

**중요: 내용 기반 최적화**
1. 텍스트 길이에 따라 박스 높이 조정 (짧으면 0.8인치, 길면 1.5-2.0인치)
2. 박스가 슬라이드를 벗어나지 않도록 (y + height ≤ 7.0)
3. 박스 간격 최소 0.15인치
4. 중요 내용은 상단 배치, 큰 폰트 사용

**JSON 구조 (반드시 이 형식으로):**
```json
{{
  "slides": [
    {{
      "slide_number": 0,
      "type": "cover",
      "elements": [
        {{"type": "gradient_bg", "x": 0, "y": 0, "width": 10, "height": 7.5, "color1": "{theme['primary']}", "color2": "{theme['secondary']}"}},
        {{"type": "accent_bar", "x": 0, "y": 0, "width": 0.12, "height": 7.5, "color": "{theme['accent']}"}},
        {{"type": "text", "x": 1.0, "y": 2.8, "width": 8.0, "height": 2.0, "text": "{project_title}", "font_size": 48, "color": "{theme['text_inverse']}", "bold": true, "align": "left"}},
        {{"type": "text", "x": 1.0, "y": 5.3, "width": 8.0, "height": 0.4, "text": "연구책임자 | {datetime.now().year}", "font_size": 15, "color": "{theme['text_inverse']}", "bold": false, "align": "left"}}
      ]
    }},
    {{
      "slide_number": 1,
      "type": "content",
      "elements": [
        {{"type": "background", "x": 0, "y": 0, "width": 10, "height": 7.5, "color": "{theme['bg_light']}"}},
        {{"type": "line", "x": 0.4, "y": 1.08, "width": 9.2, "height": 0.015, "color": "{theme['primary']}"}},
        {{"type": "text", "x": 0.4, "y": 0.25, "width": 9.2, "height": 0.75, "text": "슬라이드 제목", "font_size": 26, "color": "{theme['primary']}", "bold": true, "align": "left"}},

        {{"type": "box", "x": 1.0, "y": 1.5, "width": 8.0, "height": 1.0, "bg_color": "{theme['card_bg']}", "border_left_color": "{theme['primary']}", "border_left_width": 4, "radius": 8}},
        {{"type": "text", "x": 1.25, "y": 1.65, "width": 7.5, "height": 0.25, "text": "부제목", "font_size": 18, "color": "{theme['primary']}", "bold": true}},
        {{"type": "text", "x": 1.25, "y": 1.95, "width": 7.5, "height": 0.45, "text": "내용...", "font_size": 14, "color": "{theme['text_light']}"}},

        {{"type": "box", "x": 1.0, "y": 2.65, "width": 8.0, "height": 1.2, "bg_color": "{theme['card_bg']}", "border_left_color": "{theme['primary']}", "border_left_width": 4, "radius": 8}},
        {{"type": "text", "x": 1.25, "y": 2.8, "width": 7.5, "height": 0.25, "text": "부제목2", "font_size": 18, "color": "{theme['primary']}", "bold": true}},
        {{"type": "text", "x": 1.25, "y": 3.1, "width": 7.5, "height": 0.65, "text": "더 긴 내용...", "font_size": 14, "color": "{theme['text_light']}"}}
      ]
    }}
  ]
}}
```

**규칙:**
- 표지(slide 0) + 내용 슬라이드들
- 모든 좌표는 인치 단위 (inch)
- type: gradient_bg, background, accent_bar, line, box, text
- 텍스트는 박스 위에 별도로 배치 (박스와 겹치지 않게)
- 박스 y 위치: 1.4 이상, y+height ≤ 7.0
- 내용 길이에 따라 박스 높이 자동 조정

**완전한 JSON만 출력하세요.**"""

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                temperature=0.6,
                max_output_tokens=8000
            )
        )

        structure = json.loads(response.text)

        # 검증
        if validate_structure(structure, len(slides_data)):
            print(f"  ✓ LLM이 {len(structure['slides'])}개 슬라이드 구조 생성")
            return structure
        else:
            print(f"  ! LLM 구조 검증 실패, 폴백 사용")
            return generate_fallback_structure(slides_data, project_title, theme)

    except Exception as e:
        print(f"  [LLM 실패] {e}")
        return generate_fallback_structure(slides_data, project_title, theme)


def validate_structure(structure: Dict, expected_slides: int) -> bool:
    """구조 검증"""
    if not structure or 'slides' not in structure:
        return False
    slides = structure['slides']
    if len(slides) != expected_slides + 1:  # +1 for cover
        return False
    for slide in slides:
        if 'elements' not in slide or not slide['elements']:
            return False
    return True


# ============================================
# 폴백 구조 생성
# ============================================

def generate_fallback_structure(slides_data: List[dict], title: str, theme: Dict) -> Dict:
    """폴백 PPTX 구조"""
    structure = {"slides": []}

    # 표지
    structure["slides"].append({
        "slide_number": 0,
        "type": "cover",
        "elements": [
            {"type": "gradient_bg", "x": 0, "y": 0, "width": 10, "height": 7.5,
             "color1": theme['primary'], "color2": theme['secondary']},
            {"type": "accent_bar", "x": 0, "y": 0, "width": 0.12, "height": 7.5,
             "color": theme['accent']},
            {"type": "text", "x": 1.0, "y": 2.8, "width": 8.0, "height": 2.0,
             "text": title, "font_size": 36, "color": theme['text_inverse'],
             "bold": True, "align": "left"},
            {"type": "text", "x": 1.0, "y": 5.3, "width": 8.0, "height": 0.4,
             "text": f"연구책임자 | {datetime.now().year}", "font_size": 15,
             "color": theme['text_inverse'], "bold": False, "align": "left"}
        ]
    })

    # 내용 슬라이드들
    for idx, slide in enumerate(slides_data, 1):
        elements = [
            {"type": "background", "x": 0, "y": 0, "width": 10, "height": 7.5,
             "color": theme['bg_light']},
            {"type": "line", "x": 0.4, "y": 1.08, "width": 9.2, "height": 0.015,
             "color": theme['primary']},
            {"type": "text", "x": 0.4, "y": 0.3, "width": 9.2, "height": 0.7,
             "text": slide.get('title', ''), "font_size": 26,
             "color": theme['primary'], "bold": True, "align": "left"}
        ]

        # 박스들
        y = 1.5
        for item in slide.get('items', []):
            if y >= 7.0:
                break

            subtitle = item.get('subtitle', '')
            content = item.get('content', '')
            total_text = subtitle + content

            # 텍스트 길이에 따른 높이
            box_h = 0.8 if len(total_text) < 100 else (1.1 if len(total_text) < 300 else 1.5)
            box_h = min(box_h, 7.0 - y - 0.1)

            if box_h < 0.4:
                break

            # 박스 배경
            elements.append({
                "type": "box", "x": 1.0, "y": y, "width": 8.0, "height": box_h,
                "bg_color": theme['card_bg'], "border_left_color": theme['primary'],
                "border_left_width": 4, "radius": 8
            })

            # 부제목
            if subtitle:
                elements.append({
                    "type": "text", "x": 1.25, "y": y + 0.15, "width": 7.5, "height": 0.25,
                    "text": subtitle, "font_size": 18, "color": theme['primary'],
                    "bold": True
                })

            # 내용
            if content:
                content_y = y + 0.15 + (0.3 if subtitle else 0)
                content_h = box_h - 0.3 - (0.3 if subtitle else 0)
                elements.append({
                    "type": "text", "x": 1.25, "y": content_y, "width": 7.5, "height": content_h,
                    "text": content[:500], "font_size": 14, "color": theme['text_light'],
                    "bold": False
                })

            y += box_h + 0.15

        structure["slides"].append({
            "slide_number": idx,
            "type": "content",
            "elements": elements
        })

    return structure


# ============================================
# HTML 생성 (미리보기용)
# ============================================

def generate_preview_html(structure: Dict, theme: Dict) -> str:
    """JSON 구조를 HTML 미리보기로 변환"""
    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: 'Segoe UI', 'Malgun Gothic', sans-serif; background: #e0e0e0; padding: 20px; }}
  .slide {{
    width: 1000px;
    height: 750px;
    position: relative;
    margin: 20px auto;
    box-shadow: 0 8px 16px rgba(0,0,0,0.15);
    page-break-after: always;
    overflow: hidden;
    background: white;
  }}
  .element {{
    position: absolute;
    box-sizing: border-box;
  }}
</style>
</head>
<body>
"""

    for slide_data in structure.get('slides', []):
        slide_num = slide_data.get('slide_number', 0)
        html += f'  <section class="slide" data-slide-number="{slide_num}">\n'

        for elem in slide_data.get('elements', []):
            elem_type = elem.get('type')
            x_px = int(elem['x'] * 100)
            y_px = int(elem['y'] * 100)
            w_px = int(elem['width'] * 100)
            h_px = int(elem['height'] * 100)

            if elem_type == 'gradient_bg':
                color1 = elem.get('color1', '#065A82')
                color2 = elem.get('color2', '#1C7293')
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; background:linear-gradient(135deg, {color1}, {color2});"></div>\n'

            elif elem_type == 'background':
                color = elem.get('color', '#F8F9FB')
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; background:{color};"></div>\n'

            elif elem_type == 'accent_bar':
                color = elem.get('color', '#21295C')
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; background:{color};"></div>\n'

            elif elem_type == 'line':
                color = elem.get('color', '#1C7293')
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; background:{color};"></div>\n'

            elif elem_type == 'box':
                bg_color = elem.get('bg_color', '#F0F5FA')
                border_color = elem.get('border_left_color', '#065A82')
                border_w = elem.get('border_left_width', 4)
                radius = elem.get('radius', 8)
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; background:{bg_color}; border-left:{border_w}px solid {border_color}; border-radius:{radius}px; box-shadow:0 2px 4px rgba(0,0,0,0.1);"></div>\n'

            elif elem_type == 'text':
                text = elem.get('text', '')
                font_size = elem.get('font_size', 14)
                color = elem.get('color', '#000000')
                bold = 'font-weight:bold;' if elem.get('bold', False) else ''
                align = elem.get('align', 'left')
                html += f'    <div class="element" style="left:{x_px}px; top:{y_px}px; width:{w_px}px; height:{h_px}px; font-size:{font_size}px; color:{color}; {bold} text-align:{align}; display:flex; align-items:flex-start; padding:5px;">{text}</div>\n'

        html += '  </section>\n'

    html += '</body>\n</html>'
    return html


def generate_html(slides_data: List[dict], project_title: str, theme: Dict) -> str:
    """LLM이 내용을 분석하여 최적의 레이아웃 HTML 생성"""
    try:
        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            return generate_fallback_html(slides_data, project_title, theme)

        client = genai.Client(api_key=api_key)
        slides_with_images = prepare_images(slides_data)

        # 슬라이드 내용 분석 요약
        content_analysis = []
        for idx, slide in enumerate(slides_with_images, 1):
            items = slide.get("items", [])
            total_text_length = sum(len(item.get('subtitle', '') + item.get('content', '')) for item in items)
            content_analysis.append({
                "slide_num": idx,
                "title": slide.get("title", "")[:50],
                "item_count": len(items),
                "avg_text_length": total_text_length // max(len(items), 1),
                "has_long_content": total_text_length > 500
            })

        prompt = f"""당신은 전문 프레젠테이션 디자이너입니다.
**내용을 분석하여 각 슬라이드에 최적화된 레이아웃의 HTML을 생성**하세요.

**프로젝트:** {project_title}

**슬라이드 내용 분석:**
{json.dumps(content_analysis, ensure_ascii=False, indent=2)}

**전체 데이터:**
{json.dumps(slides_with_images, ensure_ascii=False, indent=2)}

**색상 테마:**
Primary: {theme['primary']}, Secondary: {theme['secondary']}, Accent: {theme['accent']}
Text Main: {theme['text_main']}, Text Light: {theme['text_light']}, Inverse: {theme['text_inverse']}
Background: {theme['bg_light']}, Card: {theme['card_bg']}

**중요: 내용 기반 레이아웃 최적화**
1. **텍스트 길이에 따른 박스 높이 조정**
   - 짧은 내용 (100자 미만): 박스 높이 60-80px
   - 중간 내용 (100-300자): 박스 높이 90-130px
   - 긴 내용 (300자 이상): 박스 높이 140-200px

2. **내용 중요도 분석**
   - 핵심 키워드 포함 항목: 더 큰 폰트(18-20px), 상단 배치
   - 세부 설명 항목: 작은 폰트(14-16px), 하단 배치
   - 슬라이드당 최대 4-5개 항목만 표시 (공간 부족 시)

3. **절대 위치 지정**
   - **모든 텍스트 요소(h1, h3, p)에도 position:absolute + left/top을 명시**
   - 예: <h3 style="position:absolute; left:0; top:0; ...">
   - 예: <p style="position:absolute; left:0; top:30px; ...">

4. **파싱 친화적 구조**
   - 클래스명: slide-header, content-box, box-background, box-content
   - 모든 크기와 위치는 px 단위로 명시

**표지 슬라이드 구조 (반드시 포함):**
```html
<section class="slide slide-cover" data-slide-number="0" style="position:relative; width:1000px; height:750px;">
  <div style="position:absolute; left:0; top:0; width:1000px; height:750px;
       background:linear-gradient(135deg, {theme['primary']}, {theme['secondary']});"></div>
  <div style="position:absolute; left:0; top:0; width:12px; height:750px; background:{theme['accent']};"></div>
  <div style="position:absolute; left:100px; top:280px; width:800px;">
    <h1 style="position:absolute; left:0; top:0; font-size:48px; color:{theme['text_inverse']};
         font-weight:bold; line-height:1.3;">{project_title}</h1>
  </div>
  <div style="position:absolute; left:100px; top:530px; width:800px;">
    <p style="position:absolute; left:0; top:0; font-size:16px; color:{theme['text_inverse']};">
      연구책임자 | {datetime.now().year}</p>
  </div>
</section>
```

**내용 슬라이드 예시:**
```html
<section class="slide slide-content" data-slide-number="1" style="position:relative; width:1000px; height:750px;">
  <div style="position:absolute; left:0; top:0; width:1000px; height:750px; background:{theme['bg_light']};"></div>

  <div class="slide-header" style="position:absolute; left:50px; top:30px; width:900px; height:80px;">
    <div style="position:absolute; left:0; top:77px; width:900px; height:3px; background:{theme['secondary']};"></div>
    <h1 style="position:absolute; left:0; top:20px; font-size:28px; color:{theme['primary']}; font-weight:bold;">
      슬라이드 제목</h1>
  </div>

  <!-- 박스 1: 내용이 짧으면 높이 80px -->
  <div class="content-box" style="position:absolute; left:100px; top:150px; width:800px; height:80px;">
    <div class="box-background" style="position:absolute; left:0; top:0; width:800px; height:80px;
         background:{theme['card_bg']}; border-left:4px solid {theme['primary']}; border-radius:8px;"></div>
    <div class="box-content" style="position:absolute; left:20px; top:12px; width:760px; height:56px;">
      <h3 style="position:absolute; left:0; top:0; font-size:18px; color:{theme['primary']};
           font-weight:bold;">부제목</h3>
      <p style="position:absolute; left:0; top:26px; font-size:14px; color:{theme['text_light']};
          line-height:1.5;">짧은 내용...</p>
    </div>
  </div>

  <!-- 박스 2: 내용이 길면 높이 150px -->
  <div class="content-box" style="position:absolute; left:100px; top:245px; width:800px; height:150px;">
    <div class="box-background" style="position:absolute; left:0; top:0; width:800px; height:150px;
         background:{theme['card_bg']}; border-left:4px solid {theme['primary']}; border-radius:8px;"></div>
    <div class="box-content" style="position:absolute; left:20px; top:12px; width:760px; height:126px;">
      <h3 style="position:absolute; left:0; top:0; font-size:18px; color:{theme['primary']}; font-weight:bold;">
        긴 내용의 부제목</h3>
      <p style="position:absolute; left:0; top:26px; font-size:14px; color:{theme['text_light']};
          line-height:1.5; white-space:pre-wrap;">긴 내용 텍스트...
여러 줄에 걸쳐...</p>
    </div>
  </div>
</section>
```

**필수 규칙:**
- 모든 요소에 position:absolute + left/top (px 단위)
- h1, h3, p 태그에도 position:absolute 반드시 포함
- 박스 높이는 내용 길이에 맞게 동적 조정
- 슬라이드 하단(750px)을 넘지 않도록 배치
- 박스 간격 최소 15px 유지

완전한 HTML만 출력하세요 (<!DOCTYPE html>로 시작)."""

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(temperature=0.7, max_output_tokens=15000)
        )

        html = clean_html(response.text)

        # HTML 검증 후 반환
        if validate_html(html):
            print(f"  ✓ LLM이 내용 분석 기반 HTML 생성 완료")
            return html
        else:
            print(f"  ! LLM HTML 검증 실패, 폴백 사용")
            return generate_fallback_html(slides_data, project_title, theme)

    except Exception as e:
        print(f"  [LLM 실패] {e}")
        return generate_fallback_html(slides_data, project_title, theme)


def prepare_images(slides_data: List[dict]) -> List[dict]:
    """이미지 base64 변환"""
    result = []
    for slide in slides_data:
        slide_copy = slide.copy()
        slide_copy["items"] = []
        
        for item in slide.get("items", []):
            item_copy = item.copy()
            if item.get("image_path") and os.path.exists(item["image_path"]):
                try:
                    with open(item["image_path"], "rb") as f:
                        img_data = base64.b64encode(f.read()).decode()
                        ext = item["image_path"].split('.')[-1].lower()
                        mime = f"image/{ext if ext in ['png','jpg','jpeg','gif'] else 'png'}"
                        item_copy["image_base64"] = f"data:{mime};base64,{img_data}"
                except:
                    pass
            slide_copy["items"].append(item_copy)
        result.append(slide_copy)
    
    return result


def clean_html(text: str) -> str:
    """HTML 추출"""
    text = re.sub(r'^```html\s*\n', '', text.strip(), flags=re.MULTILINE)
    text = re.sub(r'\n```$', '', text)
    
    match = re.search(r'(<!DOCTYPE[^>]*>.*)', text, re.DOTALL | re.IGNORECASE)
    if match:
        return match.group(1).strip()
    
    match = re.search(r'(<html[^>]*>.*)', text, re.DOTALL | re.IGNORECASE)
    if match:
        return match.group(1).strip()
    
    return text.strip()


def validate_html(html: str) -> bool:
    """HTML 검증"""
    if not html:
        return False
    h = html.lower()
    return '<html' in h and '<body' in h and 'class="slide"' in h


def calculate_text_height(text: str, font_size: int, width_px: int, line_height: float = 1.5) -> int:
    """텍스트 높이 계산 (대략적)"""
    if not text:
        return 0
    # 대략적인 계산: 한글 1자 ≈ font_size * 0.8px, 영문/숫자 ≈ font_size * 0.5px
    chars_per_line = width_px / (font_size * 0.7)  # 평균
    lines = max(1, len(text) / chars_per_line)
    return int(lines * font_size * line_height) + 10  # 여백 추가


def generate_fallback_html(slides_data: List[dict], title: str, theme: Dict) -> str:
    """폴백 HTML - 표지 포함, 텍스트 길이에 따른 동적 높이"""
    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: 'Segoe UI', 'Malgun Gothic', sans-serif; background: #e0e0e0; }}
  .slide {{
    width: 1000px;
    height: 750px;
    position: relative;
    margin: 20px auto;
    box-shadow: 0 8px 16px rgba(0,0,0,0.15);
    page-break-after: always;
    overflow: hidden;
  }}
</style>
</head>
<body>
  <!-- 표지 슬라이드 -->
  <section class="slide slide-cover" data-slide-number="0">
    <!-- 그라데이션 배경 -->
    <div style="position: absolute; left: 0; top: 0; width: 1000px; height: 750px;
         background: linear-gradient(135deg, {theme['primary']}, {theme['secondary']});"></div>

    <!-- 좌측 악센트 바 -->
    <div style="position: absolute; left: 0; top: 0; width: 12px; height: 750px;
         background: {theme['accent']};"></div>

    <!-- 제목 -->
    <div style="position: absolute; left: 100px; top: 280px; width: 800px;">
      <h1 style="position: absolute; left: 0; top: 0; font-size: 48px; color: {theme['text_inverse']}; font-weight: bold; line-height: 1.3;">
        {title}
      </h1>
    </div>

    <!-- 부제 -->
    <div style="position: absolute; left: 100px; top: 530px; width: 800px;">
      <p style="position: absolute; left: 0; top: 0; font-size: 16px; color: {theme['text_inverse']}; opacity: 0.9;">
        연구책임자 | {datetime.now().year}
      </p>
    </div>
  </section>
"""

    # 내용 슬라이드들
    for idx, slide in enumerate(slides_data, 1):
        html += f"""
  <!-- 슬라이드 {idx} -->
  <section class="slide slide-content" data-slide-number="{idx}">
    <!-- 배경 -->
    <div style="position: absolute; left: 0; top: 0; width: 1000px; height: 750px;
         background: {theme['bg_light']};"></div>

    <!-- 헤더 영역 -->
    <div class="slide-header" style="position: absolute; left: 50px; top: 30px; width: 900px; height: 80px;">
      <!-- 헤더 하단 라인 -->
      <div style="position: absolute; left: 0; top: 77px; width: 900px; height: 3px;
           background: {theme['secondary']};"></div>

      <!-- 제목 텍스트 -->
      <h1 style="position: absolute; left: 0; top: 20px; font-size: 28px;
           color: {theme['primary']}; font-weight: bold;">
        {slide.get('title', '')}
      </h1>
    </div>
"""

        # 콘텐츠 박스들 - 동적 높이 계산
        y = 150
        max_y = 720  # 슬라이드 하단 여백
        items = slide.get('items', [])

        for item_idx, item in enumerate(items):
            if y >= max_y:
                break

            subtitle = item.get('subtitle', '')
            content = item.get('content', '')

            # 텍스트 높이 계산
            subtitle_h = calculate_text_height(subtitle, 18, 760) if subtitle else 0
            content_h = calculate_text_height(content, 14, 760) if content else 0

            # 박스 높이 = 제목 높이 + 내용 높이 + 패딩
            box_height = max(60, subtitle_h + content_h + 30)

            # 슬라이드를 벗어나면 높이 조정
            if y + box_height > max_y:
                box_height = max_y - y
                if box_height < 40:
                    break

            html += f"""
    <!-- 콘텐츠 박스 {item_idx + 1} -->
    <div class="content-box" style="position: absolute; left: 100px; top: {y}px; width: 800px; height: {box_height}px;">
      <!-- 박스 배경 -->
      <div class="box-background" style="position: absolute; left: 0; top: 0; width: 800px; height: {box_height}px;
           background: {theme['card_bg']}; border-left: 4px solid {theme['primary']};
           border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);"></div>

      <!-- 박스 내용 -->
      <div class="box-content" style="position: absolute; left: 20px; top: 12px; width: 760px; height: {box_height - 24}px;">
        <h3 style="position: absolute; left: 0; top: 0; font-size: 18px; color: {theme['primary']}; font-weight: bold; margin-bottom: 6px;">
          {subtitle}
        </h3>
        <p style="position: absolute; left: 0; top: {subtitle_h + 8}px; font-size: 14px; color: {theme['text_light']}; line-height: 1.5; white-space: pre-wrap;">
          {content}
        </p>
      </div>
    </div>
"""
            y += box_height + 15  # 박스 간격

        html += "  </section>\n"

    html += "</body>\n</html>"
    return html


# ============================================
# HTML 파싱
# ============================================

def parse_html(html: str) -> List[Dict]:
    """HTML 파싱 - 배경/테두리 포함"""
    soup = BeautifulSoup(html, 'html.parser')
    slides = []
    
    for section in soup.find_all('section', class_='slide'):
        slide = {
            "slide_number": int(section.get('data-slide-number', 0)),
            "elements": [],
            "is_cover": 'slide-cover' in section.get('class', [])
        }
        
        # 배경 div 찾기
        for bg_div in section.find_all('div', recursive=False):
            style = parse_style(bg_div.get('style', ''))
            
            # 배경인지 판단 (position:absolute + 전체 크기)
            if (style.get('position') == 'absolute' and 
                'background' in bg_div.get('style', '') and
                'width' in style and 'height' in style):
                
                width_val = parse_length(style.get('width', '1000px'), 10.0)
                height_val = parse_length(style.get('height', '750px'), 7.5)
                
                # 슬라이드 크기 배경이면 추가
                if width_val >= 9.5 and height_val >= 7.0:
                    bg_color = extract_color_from_style(style.get('background', ''))
                    if bg_color:
                        slide["bg_color"] = bg_color
        
        # 헤더 영역 파싱
        header = section.find(class_='slide-header')
        if header:
            # 헤더 내부의 h1 찾기
            h1 = header.find('h1')
            if h1:
                elem = extract_text_element(h1, 'title', header)
                if elem:
                    slide["elements"].append(elem)
            
            # 헤더 라인 (border-bottom 또는 별도 div)
            for line_div in header.find_all('div'):
                line_style = parse_style(line_div.get('style', ''))
                if line_style.get('height') in ['3px', '2px', '4px']:  # 라인으로 판단
                    elem = create_line_element(line_div)
                    if elem:
                        slide["elements"].append(elem)
        
        # 콘텐츠 박스들 파싱
        for box in section.find_all(class_='content-box'):
            # 박스 배경
            bg = box.find(class_='box-background')
            if bg:
                elem = create_box_background(bg)
                if elem:
                    slide["elements"].append(elem)
            
            # 박스 내용
            content_div = box.find(class_='box-content')
            if content_div:
                # 부제목
                h3 = content_div.find('h3')
                if h3:
                    elem = extract_text_element(h3, 'subtitle', box)
                    if elem:
                        slide["elements"].append(elem)
                
                # 내용
                p = content_div.find('p')
                if p:
                    elem = extract_text_element(p, 'content', box)
                    if elem:
                        slide["elements"].append(elem)
        
        # 표지 슬라이드 특수 처리
        if slide["is_cover"]:
            # 악센트 바
            for div in section.find_all('div', recursive=False):
                style = parse_style(div.get('style', ''))
                width = style.get('width', '')
                if width in ['12px', '10px', '15px']:  # 좁은 바
                    elem = create_accent_bar(div)
                    if elem:
                        slide["elements"].append(elem)
            
            # 제목/부제 div
            for div in section.find_all('div', recursive=False):
                h1 = div.find('h1')
                p = div.find('p')
                
                if h1:
                    elem = extract_text_element(h1, 'cover_title', div)
                    if elem:
                        slide["elements"].append(elem)
                
                if p:
                    elem = extract_text_element(p, 'cover_subtitle', div)
                    if elem:
                        slide["elements"].append(elem)
        
        # 이미지
        for img in section.find_all('img'):
            elem = extract_element(img, 'image')
            if elem:
                slide["elements"].append(elem)
        
        slides.append(slide)
    
    return slides


def extract_text_element(elem, elem_type: str, parent) -> Optional[Dict]:
    """텍스트 요소 추출 - 모든 부모 위치 누적"""
    try:
        # 요소 자체 위치
        elem_style = parse_style(elem.get('style', ''))

        x = parse_length(elem_style.get('left', '0'), 10.0)
        y = parse_length(elem_style.get('top', '0'), 7.5)

        # 부모 체인을 따라 올라가며 모든 absolute 위치 누적
        current = elem.parent
        while current and current.name != 'section':
            current_style = parse_style(current.get('style', ''))
            if current_style.get('position') == 'absolute':
                x += parse_length(current_style.get('left', '0'), 10.0)
                y += parse_length(current_style.get('top', '0'), 7.5)
            current = current.parent
        
        # 크기
        w = parse_length(elem_style.get('width', '5in'), 10.0)
        h = parse_length(elem_style.get('height', '1in'), 7.5)
        
        # 폰트
        font_size = parse_length(elem_style.get('font-size', '14px'), 72.0) * 0.75  # px to pt
        color = elem_style.get('color', '#000000')
        
        return {
            "type": elem_type,
            "x": x,
            "y": y,
            "width": w if w > 0.5 else 5.0,
            "height": h if h > 0.3 else 0.5,
            "text": elem.get_text().strip(),
            "font_size": int(font_size) if font_size > 8 else 14,
            "color": color,
            "bold": 'bold' in elem_style.get('font-weight', '')
        }
    except:
        return None


def create_box_background(elem) -> Optional[Dict]:
    """박스 배경 요소 생성"""
    try:
        style = parse_style(elem.get('style', ''))

        x = parse_length(style.get('left', '0'), 10.0)
        y = parse_length(style.get('top', '0'), 7.5)
        w = parse_length(style.get('width', '5in'), 10.0)
        h = parse_length(style.get('height', '1in'), 7.5)

        # 부모 체인을 따라 올라가며 모든 absolute 위치 누적
        current = elem.parent
        while current and current.name != 'section':
            current_style = parse_style(current.get('style', ''))
            if current_style.get('position') == 'absolute':
                x += parse_length(current_style.get('left', '0'), 10.0)
                y += parse_length(current_style.get('top', '0'), 7.5)
            current = current.parent
        
        bg_color = extract_color_from_style(style.get('background', ''))
        border_color = extract_color_from_style(style.get('border-left', ''))
        
        return {
            "type": "box_bg",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "bg_color": bg_color,
            "border_color": border_color
        }
    except:
        return None


def create_line_element(elem) -> Optional[Dict]:
    """라인 요소 생성"""
    try:
        style = parse_style(elem.get('style', ''))

        x = parse_length(style.get('left', '0'), 10.0)
        # bottom이 지정된 경우 처리 (부모 높이에서 역산)
        if 'bottom' in style:
            # bottom: 0 → 부모의 하단, 즉 부모의 top + height
            parent = elem.parent
            if parent:
                parent_style = parse_style(parent.get('style', ''))
                parent_y = parse_length(parent_style.get('top', '0'), 7.5)
                parent_h = parse_length(parent_style.get('height', '80px'), 7.5)
                y = parent_y + parent_h - parse_length(style.get('bottom', '0'), 7.5)
            else:
                y = 7.5 - parse_length(style.get('bottom', '0'), 7.5)
        else:
            y = parse_length(style.get('top', '0'), 7.5)

        w = parse_length(style.get('width', '5in'), 10.0)
        h = parse_length(style.get('height', '3px'), 7.5)

        # 부모 위치 누적 (bottom으로 계산한 경우 제외)
        if 'bottom' not in style:
            parent = elem.parent
            if parent:
                parent_style = parse_style(parent.get('style', ''))
                if parent_style.get('position') == 'absolute':
                    parent_x = parse_length(parent_style.get('left', '0'), 10.0)
                    parent_y = parse_length(parent_style.get('top', '0'), 7.5)
                    x += parent_x
                    y += parent_y
        else:
            # bottom 계산 시 left만 누적
            parent = elem.parent
            if parent:
                parent_style = parse_style(parent.get('style', ''))
                if parent_style.get('position') == 'absolute':
                    parent_x = parse_length(parent_style.get('left', '0'), 10.0)
                    x += parent_x
        
        color = extract_color_from_style(style.get('background', ''))
        
        return {
            "type": "line",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "color": color
        }
    except:
        return None


def create_accent_bar(elem) -> Optional[Dict]:
    """악센트 바 생성"""
    try:
        style = parse_style(elem.get('style', ''))

        x = parse_length(style.get('left', '0'), 10.0)
        y = parse_length(style.get('top', '0'), 7.5)
        w = parse_length(style.get('width', '12px'), 10.0)
        h = parse_length(style.get('height', '750px'), 7.5)

        # 부모 체인을 따라 올라가며 모든 absolute 위치 누적
        current = elem.parent
        while current and current.name != 'section':
            current_style = parse_style(current.get('style', ''))
            if current_style.get('position') == 'absolute':
                x += parse_length(current_style.get('left', '0'), 10.0)
                y += parse_length(current_style.get('top', '0'), 7.5)
            current = current.parent

        color = extract_color_from_style(style.get('background', ''))

        return {
            "type": "accent_bar",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "color": color
        }
    except:
        return None


def extract_color_from_style(style_value: str) -> Optional[str]:
    """스타일 값에서 색상 추출"""
    if not style_value:
        return None
    
    # #RRGGBB
    match = re.search(r'#[0-9a-fA-F]{6}', style_value)
    if match:
        return match.group(0)
    
    # rgb(...)
    match = re.search(r'rgb\([^)]+\)', style_value)
    if match:
        return match.group(0)
    
    # gradient에서 첫 색상
    if 'gradient' in style_value:
        match = re.search(r'#[0-9a-fA-F]{6}', style_value)
        if match:
            return match.group(0)
    
    return None


def extract_element(elem, elem_type: str) -> Optional[Dict]:
    """요소 데이터 추출"""
    try:
        # data-* 속성 우선
        x = float(elem.get('data-x', 1.0))
        y = float(elem.get('data-y', 1.0))
        w = float(elem.get('data-width', 5.0))
        h = float(elem.get('data-height', 1.0))
        
        # 인라인 스타일 폴백
        style = parse_style(elem.get('style', ''))
        if 'left' in style and not elem.get('data-x'):
            x = parse_length(style['left'], 10.0)
        if 'top' in style and not elem.get('data-y'):
            y = parse_length(style['top'], 7.5)
        if 'width' in style and not elem.get('data-width'):
            w = parse_length(style['width'], 10.0)
        
        base = {"x": x, "y": y, "width": w, "height": h}
        
        if elem_type == 'title':
            return {**base, "type": "title", "text": elem.get_text().strip(),
                    "font_size": int(elem.get('data-fontsize', 28)), 
                    "color": elem.get('data-color', '#065a82'), "bold": True}
        
        elif elem_type == 'text':
            subtitle = elem.find(['h3', 'h4'])
            content = elem.find('p')
            return {**base, "type": "text",
                    "subtitle": subtitle.get_text().strip() if subtitle else "",
                    "content": content.get_text().strip() if content else elem.get_text().strip(),
                    "font_size": int(elem.get('data-fontsize', 14)),
                    "color": elem.get('data-color', '#333333'),
                    "is_highlight": 'highlight' in ' '.join(elem.get('class', []))}
        
        elif elem_type == 'image':
            return {**base, "type": "image", "src": elem.get('src', '')}
    
    except:
        return None


def parse_style(style_str: str) -> Dict:
    """CSS 파싱"""
    result = {}
    for rule in style_str.split(';'):
        if ':' in rule:
            k, v = rule.split(':', 1)
            result[k.strip().lower()] = v.strip()
    return result


def parse_length(value: str, ref: float) -> float:
    """CSS 길이 → 인치"""
    try:
        if 'px' in value:
            return float(value.replace('px', '')) / 100.0
        elif '%' in value:
            return float(value.replace('%', '')) / 100.0 * ref
        return float(value)
    except:
        return 1.0


def extract_bg_color(elem) -> Optional[str]:
    """배경색 추출"""
    style = elem.get('style', '')
    match = re.search(r'background-color:\s*([^;]+)', style)
    if match:
        return match.group(1).strip()
    match = re.search(r'background:\s*([^;]+)', style)
    if match:
        bg = match.group(1).strip()
        color = re.search(r'#[0-9a-fA-F]{6}', bg)
        if color:
            return color.group(0)
    return None


# ============================================
# JSON → PPTX 직접 렌더링
# ============================================

def render_pptx_from_structure(structure: Dict, output_path: str, theme: Dict):
    """JSON 구조를 기반으로 PPTX 직접 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in structure['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for elem in slide_data['elements']:
            elem_type = elem.get('type')

            try:
                if elem_type == 'gradient_bg':
                    render_gradient_bg(slide, elem)
                elif elem_type == 'background':
                    render_background(slide, elem)
                elif elem_type == 'accent_bar':
                    render_accent_bar(slide, elem)
                elif elem_type == 'line':
                    render_line(slide, elem)
                elif elem_type == 'box':
                    render_box(slide, elem)
                elif elem_type == 'text':
                    render_text(slide, elem)
            except Exception as e:
                print(f"  ! 요소 렌더링 실패 ({elem_type}): {e}")
                continue

    prs.save(output_path)
    print(f"  ✓ PPTX 저장: {output_path}")


def render_gradient_bg(slide, elem):
    """그라데이션 배경 (단색으로 대체)"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(elem['x']), Inches(elem['y']),
        Inches(elem['width']), Inches(elem['height'])
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = parse_color(elem['color1'])
    bg.line.fill.background()
    bg.shadow.inherit = False


def render_background(slide, elem):
    """단색 배경"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(elem['x']), Inches(elem['y']),
        Inches(elem['width']), Inches(elem['height'])
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = parse_color(elem['color'])
    bg.line.fill.background()
    bg.shadow.inherit = False


def render_accent_bar(slide, elem):
    """악센트 바"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(elem['x']), Inches(elem['y']),
        Inches(elem['width']), Inches(elem['height'])
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = parse_color(elem['color'])
    bar.line.fill.background()
    bar.shadow.inherit = False


def render_line(slide, elem):
    """라인"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(elem['x']), Inches(elem['y']),
        Inches(elem['width']), Inches(elem['height'])
    )
    line.fill.solid()
    line.fill.fore_color.rgb = parse_color(elem['color'])
    line.line.fill.background()


def render_box(slide, elem):
    """박스 (배경 + 테두리)"""
    x, y = elem['x'], elem['y']
    w, h = elem['width'], elem['height']

    # 좌측 테두리 (border-left)
    if elem.get('border_left_color'):
        border_w = elem.get('border_left_width', 4) / 72.0  # pt to inch
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(border_w), Inches(h)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = parse_color(elem['border_left_color'])
        border.line.fill.background()

    # 박스 배경
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = parse_color(elem.get('bg_color', '#F0F5FA'))
    box.line.fill.background()
    box.shadow.inherit = False


def render_text(slide, elem):
    """텍스트"""
    x = max(0, min(elem['x'], 9.5))
    y = max(0, min(elem['y'], 7.0))
    w = max(0.5, min(elem['width'], 10 - x))
    h = max(0.2, min(elem['height'], 7.5 - y))

    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = 1  # 상단 정렬

    p = tf.paragraphs[0]
    p.text = elem.get('text', '')
    p.font.size = Pt(elem.get('font_size', 14))
    p.font.bold = elem.get('bold', False)
    p.font.color.rgb = parse_color(elem.get('color', '#000000'))

    # 정렬
    align = elem.get('align', 'left')
    if align == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT


# ============================================
# PPTX 변환 (기존 - 더 이상 사용 안 함)
# ============================================

def convert_to_pptx(parsed: List[Dict], output_path: str):
    """PPTX 생성 - 배경/테두리 포함"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in parsed:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 배경색
        if slide_data.get("bg_color"):
            add_bg(slide, slide_data["bg_color"])
        
        # 요소 렌더링 (z-index 순서: 배경 → 박스 → 텍스트)
        # 1. 박스 배경들 먼저
        for elem in slide_data.get("elements", []):
            if elem["type"] == "box_bg":
                add_box_background(slide, elem)
            elif elem["type"] == "accent_bar":
                add_accent_bar(slide, elem)
            elif elem["type"] == "line":
                add_line(slide, elem)
        
        # 2. 텍스트 요소들
        for elem in slide_data.get("elements", []):
            if elem["type"] in ["title", "subtitle", "content", "cover_title", "cover_subtitle"]:
                add_text(slide, elem)
            elif elem["type"] == "image":
                add_image(slide, elem)
    
    prs.save(output_path)


def add_box_background(slide, elem: Dict):
    """박스 배경 추가"""
    try:
        x = max(0, min(elem["x"], 9.5))
        y = max(0, min(elem["y"], 7.0))
        w = max(0.5, min(elem["width"], 10 - x))
        h = max(0.3, min(elem["height"], 7.5 - y))
        
        # 배경 박스
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        
        # 배경색
        bg_color = parse_color(elem.get("bg_color", "#F0F5FA"))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        
        # 테두리 (border-left 효과)
        if elem.get("border_color"):
            border_color = parse_color(elem["border_color"])
            box.line.color.rgb = border_color
            box.line.width = Pt(4)
        else:
            box.line.fill.background()
        
        # 그림자 (선택)
        box.shadow.inherit = False
        
    except Exception as e:
        print(f"  [박스 배경 실패] {e}")


def add_line(slide, elem: Dict):
    """라인 추가 (헤더 하단 라인 등)"""
    try:
        x = max(0, min(elem["x"], 9.5))
        y = max(0, min(elem["y"], 7.0))
        w = max(0.5, min(elem["width"], 10 - x))
        h = max(0.01, min(elem["height"], 0.1))  # 라인은 얇게
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        
        color = parse_color(elem.get("color", "#1C7293"))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        
    except Exception as e:
        print(f"  [라인 실패] {e}")


def add_accent_bar(slide, elem: Dict):
    """악센트 바 추가 (표지 좌측 바)"""
    try:
        x = max(0, min(elem["x"], 0.2))  # 좌측에만
        y = 0
        w = max(0.05, min(elem["width"], 0.2))
        h = 7.5
        
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        
        color = parse_color(elem.get("color", "#21295C"))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
    except Exception as e:
        print(f"  [악센트 바 실패] {e}")


def add_bg(slide, color: str):
    """배경"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    rgb = parse_color(color)
    if rgb:
        bg.fill.fore_color.rgb = rgb
    bg.line.fill.background()


def add_text(slide, elem: Dict):
    """텍스트 추가 - 모든 타입 지원"""
    try:
        x = max(0, min(elem["x"], 9.5))
        y = max(0, min(elem["y"], 7.0))
        w = max(0.5, min(elem["width"], 10 - x))
        h = max(0.3, min(elem["height"], 7.5 - y))
        
        # 텍스트 박스 생성
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1  # 중앙 정렬
        
        p = tf.paragraphs[0]
        p.text = elem.get("text", "")
        
        # 폰트 설정
        p.font.size = Pt(elem.get("font_size", 14))
        p.font.bold = elem.get("bold", False)
        p.font.color.rgb = parse_color(elem.get("color", "#000000"))
        
        # 타입별 추가 설정
        if elem["type"] in ["title", "cover_title"]:
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
        elif elem["type"] == "cover_subtitle":
            p.alignment = PP_ALIGN.LEFT
        elif elem["type"] == "subtitle":
            p.font.bold = True
        
    except Exception as e:
        print(f"  [텍스트 실패] {e}")


def add_image(slide, elem: Dict):
    """이미지"""
    try:
        src = elem.get("src", "")
        if not src.startswith("data:image"):
            return
        
        img_data = base64.b64decode(src.split(',')[1])
        temp = f"/tmp/img_{datetime.now().strftime('%Y%m%d%H%M%S%f')}.png"
        
        with open(temp, 'wb') as f:
            f.write(img_data)
        
        x = max(0, min(elem["x"], 9))
        y = max(0, min(elem["y"], 7))
        w = max(0.5, min(elem["width"], 10 - x))
        h = max(0.5, min(elem["height"], 7.5 - y))
        
        slide.shapes.add_picture(temp, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        os.remove(temp)
    except Exception as e:
        print(f"  [이미지 실패] {e}")


def parse_color(color: str) -> RGBColor:
    """색상 파싱"""
    try:
        if color.startswith('#'):
            c = color.lstrip('#')
            if len(c) == 6:
                return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
        if color.startswith('rgb'):
            m = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color)
            if m:
                return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    except:
        pass
    return RGBColor(50, 50, 50)


# ============================================
# 메인 함수
# ============================================

def generate_ppt(
    state: Dict,
    output_dir: str,
    project_title: str = None,
    project_description: str = "",
    save_html: bool = True,
    force_theme: str = None
) -> Dict:
    """
    PPT 자동 생성 (LLM이 구조 JSON 생성 → 직접 PPTX 렌더링)

    Args:
        state: {"slides": [...]}
        output_dir: 출력 디렉토리
        project_title: 프로젝트 제목
        project_description: 프로젝트 설명 (테마 선택용)
        save_html: HTML 미리보기 저장 여부
        force_theme: 테마 강제 지정 (예: "ocean_marine")

    Returns:
        {"ppt_path": "...", "html_path": "...", "theme": "...", "structure": {...}}
    """
    try:
        slides_data = state.get("slides", [])

        if not project_title:
            project_title = state.get("analyzed_json", {}).get("project_summary", {}).get("title", "연구 제안서")

        print(f"\n[PPT 생성 - 새로운 방식]")
        print(f"  제목: {project_title}")
        print(f"  슬라이드: {len(slides_data)}개")

        # 테마 선택
        if force_theme:
            theme_key = force_theme
            print(f"  [테마 지정] {force_theme}")
        else:
            theme_key = select_theme(state, project_title, project_description)

        theme_colors = COLOR_THEMES[theme_key]

        # LLM이 PPTX 구조 JSON 생성
        print(f"  [1/2] LLM이 최적 레이아웃 구조 생성...")
        structure = generate_ppt_structure(slides_data, project_title, theme_colors)

        # JSON → PPTX 직접 렌더링
        print(f"  [2/2] PPTX 렌더링...")
        os.makedirs(output_dir, exist_ok=True)
        pptx_path = os.path.join(output_dir, f"PPT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
        render_pptx_from_structure(structure, pptx_path, theme_colors)

        # HTML 미리보기 (선택)
        html_path = None
        if save_html:
            print(f"  [+] HTML 미리보기 생성...")
            html = generate_preview_html(structure, theme_colors)
            html_path = os.path.join(output_dir, f"preview_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html)
            print(f"  ✓ HTML: {html_path}")

        print(f"\n[완료] 🎉\n")
        return {
            "ppt_path": pptx_path,
            "html_path": html_path,
            "theme": theme_key,
            "structure": structure
        }

    except Exception as e:
        print(f"\n[실패] {e}")
        import traceback
        traceback.print_exc()
        return {"ppt_path": "", "html_path": None, "theme": "", "structure": None}

# ============================================
# 테스트
# ============================================

def test():
    """테스트"""
    test_slides = [
    {
        "page_number": 1,
        "section": "연구개발 목표 및 과정",
        "title": "연구개발 목표 및 추진 과정",
        "items": [
            {
                "subtitle": "연구개발 목표",
                "content": (
                    "• 예측 정확도 15% 이상 향상을 통한 세계 3위권 수준 달성\n"
                    "• 다양한 환경 변화에 대응 가능한 안정적 AI 예측 모델 구축"
                )
            },
            {
                "subtitle": "연구개발 추진 과정",
                "content": (
                    "• 데이터 수집 및 전처리: 로그·센서 데이터 정제 및 학습 데이터셋 구축\n"
                    "• 예측 모델 설계: 도메인 특성을 반영한 시계열 기반 모델 구조 설계\n"
                    "• AI 오차 보정: 예측 오차 패턴 학습을 통한 보정 모듈 결합\n"
                    "• 성능 검증 및 고도화: 시나리오 기반 검증 및 반복 학습 수행"
                )
            }
        ]
    }
]
    
    state = {"slides": test_slides}
    
    result = generate_ppt(
        state=state,
        output_dir=r"C:\big_project\modeling\data\pptx",
        project_title="차세대 해양·극지 환경 기후예측시스템 개발"
    )
    
    print(f"\n{'='*60}")
    if result["ppt_path"]:
        print(f"성공!")
        print(f"  테마: {result['theme']}")
        print(f"  HTML: {result['html_path']}")
        print(f"  PPTX: {result['ppt_path']}")
    else:
        print(f"실패")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    test()